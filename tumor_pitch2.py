# ============================================================
#  Tumor Documentation Pitch Deck – DESIGN C
#  "Dark Sidebar + Warm Gold" — Completely different from AE deck
#  Run: pip install python-pptx   then   python tumor_pitch_c.py
# ============================================================

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── DESIGN C – Warm Dark Palette ────────────────────────────
NAVY       = RGBColor(0x1E, 0x29, 0x3B)
NAVY_MID   = RGBColor(0x33, 0x41, 0x55)
NAVY_LIGHT = RGBColor(0x47, 0x55, 0x69)
GOLD       = RGBColor(0xF5, 0x9E, 0x0B)
GOLD_DARK  = RGBColor(0xD9, 0x77, 0x06)
GOLD_LT    = RGBColor(0xFF, 0xFB, 0xEB)
ROSE       = RGBColor(0xF4, 0x3F, 0x5E)
ROSE_LT    = RGBColor(0xFF, 0xF1, 0xF2)
CYAN       = RGBColor(0x06, 0xB6, 0xD4)
CYAN_LT    = RGBColor(0xEC, 0xFE, 0xFF)
MINT       = RGBColor(0x10, 0xB9, 0x81)
MINT_LT    = RGBColor(0xEC, 0xFD, 0xF5)
VIOLET     = RGBColor(0x8B, 0x5C, 0xF6)
VIOLET_LT  = RGBColor(0xF5, 0xF3, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
WARM_50    = RGBColor(0xFD, 0xFA, 0xF6)
WARM_100   = RGBColor(0xF5, 0xF0, 0xEB)
WARM_200   = RGBColor(0xE7, 0xE0, 0xD8)
SLATE_500  = RGBColor(0x64, 0x74, 0x8B)
SLATE_400  = RGBColor(0x94, 0xA3, 0xB8)
SLATE_300  = RGBColor(0xCB, 0xD5, 0xE1)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────
def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def tx(s, l, t, w, h, text, sz=14, bold=False, c=NAVY, al=PP_ALIGN.LEFT, fn="Calibri"):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = c
    p.font.name = fn; p.alignment = al
    return tb

def ml(s, l, t, w, h, lines, sz=13, c=SLATE_500, bold=False, al=PP_ALIGN.LEFT, sp=4):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.bold = bold
        p.font.color.rgb = c; p.font.name = "Calibri"
        p.alignment = al; p.space_after = Pt(sp)
    return tb

def bx(s, l, t, w, h, fill=WARM_100, border=None, bw=Pt(1)):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
         Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = bw
    else: sh.line.fill.background()
    return sh

def circ(s, l, t, sz, fill=GOLD):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL,
         Inches(l), Inches(t), Inches(sz), Inches(sz))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def sidebar(s, title, subtitle="", num=""):
    """Dark navy left sidebar — signature element of Design C"""
    bx(s, 0, 0, 3.3, 7.5, fill=NAVY)
    # Gold accent line
    bx(s, 3.3, 0, 0.06, 7.5, fill=GOLD)
    if num:
        tx(s, 0.4, 0.5, 2.5, 0.7, num, sz=52, bold=True, c=GOLD_DARK, al=PP_ALIGN.LEFT)
    tx(s, 0.4, 1.4 if num else 0.8, 2.6, 1.5, title,
       sz=22, bold=True, c=WHITE, al=PP_ALIGN.LEFT)
    if subtitle:
        tx(s, 0.4, 2.6 if num else 2.0, 2.6, 1.0, subtitle,
           sz=12, c=SLATE_400, al=PP_ALIGN.LEFT)
    # footer in sidebar
    tx(s, 0.4, 6.8, 2.5, 0.4, "EKFZ · TU Dresden",
       sz=8, c=NAVY_LIGHT)

def snum(s, n):
    tx(s, 12.4, 7.05, 0.7, 0.3, f"{n}/10", sz=9, c=SLATE_300, al=PP_ALIGN.RIGHT)

# Card with colored top bar
def topbar_card(s, l, t, w, h, bar_color, title, lines, icon=""):
    bx(s, l, t, w, h, fill=WHITE, border=WARM_200, bw=Pt(1))
    bx(s, l, t, w, 0.07, fill=bar_color)
    y = t + 0.18
    if icon:
        tx(s, l + 0.15, y, 0.4, 0.3, icon, sz=18)
        tx(s, l + 0.55, y, w - 0.8, 0.3, title, sz=13, bold=True, c=NAVY)
    else:
        tx(s, l + 0.15, y, w - 0.3, 0.3, title, sz=13, bold=True, c=NAVY)
    ml(s, l + 0.15, y + 0.35, w - 0.3, h - 0.7, lines, sz=11, c=SLATE_500)


# ═══════════════════════════════════════════════════════════
#  SLIDE 1 – TITLE (Full dark background — centered)
# ═══════════════════════════════════════════════════════════
s = blank()

# Full navy background
bx(s, 0, 0, 13.333, 7.5, fill=NAVY)

# Decorative gold circles
c1 = circ(s, -1.5, -1.5, 4.0, fill=NAVY_MID)
c2 = circ(s, 11.0, 5.5, 3.5, fill=NAVY_MID)

# Gold top line
bx(s, 0, 0, 13.333, 0.06, fill=GOLD)

# Event tag
bx(s, 4.2, 0.8, 5.0, 0.42, fill=NAVY_MID, border=GOLD, bw=Pt(1))
tx(s, 4.2, 0.82, 5.0, 0.38,
   "CLINICUM DIGITALE  ·  SPRING 2026  ·  PROJECT X",
   sz=10, bold=True, c=GOLD, al=PP_ALIGN.CENTER)

# Main title
tx(s, 1.0, 1.8, 11.3, 0.6,
   "LLM-Powered Automated",
   sz=40, bold=True, c=WHITE, al=PP_ALIGN.CENTER)
tx(s, 1.0, 2.55, 11.3, 0.6,
   "Tumor Documentation",
   sz=40, bold=True, c=GOLD, al=PP_ALIGN.CENTER)
tx(s, 1.0, 3.25, 11.3, 0.5,
   "from Clinical Text Documents",
   sz=24, bold=False, c=SLATE_400, al=PP_ALIGN.CENTER)

# Subtitle
tx(s, 2.0, 4.2, 9.3, 0.7,
   "Build a web tool that extracts structured oncological information\n"
   "from clinical documents using LLM-based RAG and ICD-O coding.",
   sz=15, c=SLATE_400, al=PP_ALIGN.CENTER)

# Three feature pills
pills = [("📄 Document Upload", GOLD), ("🧠 LLM Pipeline", CYAN), ("📊 ICD-O Coding", MINT)]
px = 3.5
for text, clr in pills:
    bx(s, px, 5.1, 2.0, 0.38, fill=NAVY_MID, border=clr, bw=Pt(1.5))
    tx(s, px, 5.12, 2.0, 0.35, text, sz=10, bold=True, c=clr, al=PP_ALIGN.CENTER)
    px += 2.2

# Supervisors
bx(s, 3.5, 5.85, 6.3, 0.85, fill=NAVY_MID, border=NAVY_LIGHT, bw=Pt(1))
# Avatar circles
iw = circ(s, 4.5, 5.98, 0.5, fill=GOLD)
iw.text_frame.paragraphs[0].text = "IW"
iw.text_frame.paragraphs[0].font.size = Pt(13); iw.text_frame.paragraphs[0].font.bold = True
iw.text_frame.paragraphs[0].font.color.rgb = NAVY
iw.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
iw.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

nd = circ(s, 5.1, 5.98, 0.5, fill=ROSE)
nd.text_frame.paragraphs[0].text = "ND"
nd.text_frame.paragraphs[0].font.size = Pt(13); nd.text_frame.paragraphs[0].font.bold = True
nd.text_frame.paragraphs[0].font.color.rgb = WHITE
nd.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
nd.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

tx(s, 5.8, 5.95, 3.8, 0.3,
   "Isabella Catharina Wiest & Naghme Dashti", sz=13, bold=True, c=WHITE)
tx(s, 5.8, 6.25, 3.8, 0.3,
   "EKFZ Digital Health · TU Dresden", sz=10, c=SLATE_400)

# Logo placeholders
for i, lb in enumerate(["EKFZ", "TUD", "CGC"]):
    lx = 10.5 + i * 1.0
    bx(s, lx, 6.6, 0.8, 0.55, fill=NAVY_MID, border=NAVY_LIGHT, bw=Pt(1))
    tx(s, lx, 6.65, 0.8, 0.45, lb, sz=8, c=SLATE_400, al=PP_ALIGN.CENTER)

snum(s, 1)


# ═══════════════════════════════════════════════════════════
#  SLIDE 2 – THE PROBLEM (Sidebar + cards right)
# ═══════════════════════════════════════════════════════════
s = blank()
sidebar(s, "The\nDocumentation\nBurden in\nOncology", "Why structured tumor\ndocumentation needs\nautomation", "01")

# Right content area
problems = [
    ("🏥", "Massive Unstructured Data", GOLD,
     ["Physician letters, pathology reports,",
      "follow-up notes — all as free text"]),
    ("⏱️", "Manual Entry Bottleneck", ROSE,
     ["30+ min per patient to transfer",
      "tumor info into registries"]),
    ("🔗", "Fragmented Records", CYAN,
     ["Tumor data scattered across",
      "multiple documents per patient"]),
    ("⚠️", "Incomplete Data Affects Care", VIOLET,
     ["Missing data impacts therapy",
      "decisions & tumor boards"]),
]

x, y = 3.7, 0.3
for i, (icon, title, clr, desc) in enumerate(problems):
    topbar_card(s, x, y, 4.55, 1.25, clr, title, desc, icon)
    if i % 2 == 0:
        x += 4.85
    else:
        x = 3.7; y += 1.45

# Before vs After
bx(s, 3.7, 3.3, 4.55, 1.35, fill=ROSE_LT, border=ROSE, bw=Pt(1.5))
tx(s, 3.9, 3.38, 4.2, 0.25, "❌  TODAY", sz=13, bold=True, c=ROSE)
tx(s, 3.9, 3.65, 4.2, 0.25,
   "📄 Read Docs → ✍️ Manual Extraction → 📋 Registry",
   sz=11, c=NAVY)
tx(s, 3.9, 3.95, 4.2, 0.25,
   "⏳ 30+ min/patient · Error-prone · Incomplete",
   sz=10, bold=True, c=ROSE)
ml(s, 3.9, 4.2, 4.2, 0.3,
   ["77 OCR docs · 50 patients · 63 GT tumor labels"],
   sz=10, c=SLATE_500)

bx(s, 8.55, 3.3, 4.55, 1.35, fill=MINT_LT, border=MINT, bw=Pt(1.5))
tx(s, 8.75, 3.38, 4.2, 0.25, "✅  OUR VISION", sz=13, bold=True, c=MINT)
tx(s, 8.75, 3.65, 4.2, 0.25,
   "📄 Upload Docs → 🤖 LLM Pipeline → 📊 Structured",
   sz=11, c=NAVY)
tx(s, 8.75, 3.95, 4.2, 0.25,
   "⚡ Automated · ICD-O structured · Complete",
   sz=10, bold=True, c=MINT)
ml(s, 8.75, 4.2, 4.2, 0.3,
   ["4,241 ICD-O dictionary terms available"],
   sz=10, c=SLATE_500)

# Bottom stat strip
bx(s, 3.7, 4.95, 9.3, 0.8, fill=NAVY)
stats = [("77", "OCR Documents"), ("50", "Patients"),
         ("63", "GT Labels"), ("4,241", "ICD-O Terms")]
sx = 4.2
for val, lab in stats:
    tx(s, sx, 5.0, 1.8, 0.3, val, sz=22, bold=True, c=GOLD, al=PP_ALIGN.CENTER)
    tx(s, sx, 5.35, 1.8, 0.25, lab, sz=9, c=SLATE_400, al=PP_ALIGN.CENTER)
    sx += 2.2

snum(s, 2)


# ═══════════════════════════════════════════════════════════
#  SLIDE 3 – ICD-O (Sidebar + two big panels right)
# ═══════════════════════════════════════════════════════════
s = blank()
sidebar(s, "ICD-O:\nThe Oncology\nCoding\nStandard",
        "Dual-axis classification:\nWHERE + WHAT", "02")

# Topography panel (dark header)
bx(s, 3.7, 0.3, 4.55, 3.2, fill=WHITE, border=WARM_200, bw=Pt(1))
bx(s, 3.7, 0.3, 4.55, 0.55, fill=CYAN)
tx(s, 3.9, 0.35, 4.2, 0.45, "📍  TOPOGRAPHY — Where?",
   sz=15, bold=True, c=WHITE)
ml(s, 3.9, 0.95, 4.2, 2.4, [
    "1,338 terms in ICD-O-3 dictionary",
    "Format: C00–C80",
    "",
    "Examples:",
    "  C73.9 → Thyroid gland",
    "  C61   → Prostate",
    "  C50.9 → Breast",
    "  C18.0 → Cecum",
], sz=12, c=SLATE_500)

# Morphology panel (dark header)
bx(s, 8.55, 0.3, 4.55, 3.2, fill=WHITE, border=WARM_200, bw=Pt(1))
bx(s, 8.55, 0.3, 4.55, 0.55, fill=ROSE)
tx(s, 8.75, 0.35, 4.2, 0.45, "🔬  MORPHOLOGY — What?",
   sz=15, bold=True, c=WHITE)
ml(s, 8.75, 0.95, 4.2, 2.4, [
    "2,903 terms in ICD-O-3 dictionary",
    "Format: XXXX/B (histology/behavior)",
    "",
    "Behavior codes:",
    "  /0 = Benign",
    "  /2 = In situ",
    "  /3 = Malignant, primary",
    "  /6 = Malignant, metastatic",
], sz=12, c=SLATE_500)

# Document types — horizontal strip
tx(s, 3.7, 3.75, 6, 0.3, "Clinical Document Types",
   sz=16, bold=True, c=NAVY)

docs = [
    ("📝", "Arztbrief", "Physician letter", GOLD),
    ("🔬", "Pathologiebefund", "Histology report", ROSE),
    ("📊", "Verlaufsdoku", "Follow-up notes", CYAN),
    ("🖼", "Befundbericht", "Imaging reports", VIOLET),
]
dx = 3.7
for icon, title, desc, clr in docs:
    bx(s, dx, 4.15, 2.2, 1.0, fill=WHITE, border=WARM_200, bw=Pt(1))
    bx(s, dx, 4.15, 2.2, 0.06, fill=clr)
    tx(s, dx + 0.1, 4.3, 0.35, 0.3, icon, sz=18)
    tx(s, dx + 0.5, 4.3, 1.6, 0.25, title, sz=11, bold=True, c=NAVY)
    tx(s, dx + 0.5, 4.58, 1.6, 0.25, desc, sz=10, c=SLATE_500)
    # Colored bottom accent
    bx(s, dx + 0.1, 4.9, 2.0, 0.15, fill=clr)
    tx(s, dx + 0.1, 4.9, 2.0, 0.15, "document-specific prompt",
       sz=7, c=WHITE, al=PP_ALIGN.CENTER)
    dx += 2.35

# Key insight box
bx(s, 3.7, 5.45, 9.3, 0.55, fill=GOLD_LT, border=GOLD, bw=Pt(1))
tx(s, 3.9, 5.5, 8.9, 0.45,
   "💡  Each document type triggers a different prompt in the LLM pipeline → "
   "type assignment is a critical user task!",
   sz=13, bold=True, c=GOLD_DARK)

snum(s, 3)


# ═══════════════════════════════════════════════════════════
#  SLIDE 4 – PIPELINE (Sidebar + Vertical waterfall right)
# ═══════════════════════════════════════════════════════════
s = blank()
sidebar(s, "Multi-Stage\nLLM Pipeline",
        "Four stages from raw\ndocuments to structured\nICD-O codes", "03")

# Vertical waterfall — stages flow top to bottom with offset
stages = [
    (1, "Document Ingestion", "📄", GOLD,
     "OCR clinical docs → merge per patient → clean & chunk"),
    (2, "LLM Summarization", "🧠", ROSE,
     "GPT-OSS-120B → map-reduce → 7–10 sentence tumor summary"),
    (3, "RAG Retrieval", "📐", CYAN,
     "BGE-M3 embeddings → cosine similarity → Top-K Morph & Topo candidates"),
    (4, "ICD-O Code Assignment", "🏷", MINT,
     "LLM reads summary + candidates → structured JSON → Morph + Topo + Reason"),
]

y = 0.3
for num, title, icon, clr, desc in stages:
    # Alternating indent
    indent = 3.7 if num % 2 == 1 else 5.5
    w = 7.5 if num % 2 == 1 else 7.5

    bx(s, indent, y, w, 1.0, fill=WHITE, border=clr, bw=Pt(2))

    # Number circle
    nc = circ(s, indent + 0.15, y + 0.15, 0.5, fill=clr)
    nc.text_frame.paragraphs[0].text = str(num)
    nc.text_frame.paragraphs[0].font.size = Pt(16)
    nc.text_frame.paragraphs[0].font.bold = True
    nc.text_frame.paragraphs[0].font.color.rgb = WHITE
    nc.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    nc.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    tx(s, indent + 0.75, y + 0.1, 0.4, 0.35, icon, sz=20)
    tx(s, indent + 1.2, y + 0.13, w - 1.5, 0.3, title,
       sz=15, bold=True, c=clr)
    tx(s, indent + 1.2, y + 0.5, w - 1.5, 0.4, desc,
       sz=12, c=SLATE_500)

    # Connecting arrow
    if num < 4:
        ax = indent + 0.38
        tx(s, ax, y + 1.0, 0.5, 0.3, "↓", sz=18, bold=True, c=SLATE_300)

    y += 1.3

# Bottom banner
bx(s, 3.7, 5.65, 9.3, 0.65, fill=NAVY)
bx(s, 3.7, 5.65, 0.08, 0.65, fill=GOLD)
tx(s, 4.0, 5.72, 8.8, 0.5,
   "💡  The entire pipeline exists as a FastAPI backend — "
   "you build the web interface!",
   sz=14, bold=True, c=GOLD, al=PP_ALIGN.CENTER)

# Tech tags
tags = [("🧠 GPT-OSS-120B", ROSE), ("📐 BGE-M3", CYAN),
        ("📖 ICD-O-3", MINT), ("⚡ FastAPI", GOLD)]
tx_x = 3.7
for text, clr in tags:
    bx(s, tx_x, 6.5, 2.2, 0.35, fill=WARM_50, border=clr, bw=Pt(1))
    tx(s, tx_x, 6.52, 2.2, 0.3, text, sz=10, bold=True, c=clr, al=PP_ALIGN.CENTER)
    tx_x += 2.4

snum(s, 4)


# ═══════════════════════════════════════════════════════════
#  SLIDE 5 – RESEARCH RESULTS (Sidebar + metric cards)
# ═══════════════════════════════════════════════════════════
s = blank()
sidebar(s, "Research\nFindings",
        "Baseline evaluation on\n50 patients with\nANY-match strategy", "04")

# Four colored metric cards in 2x2
findings = [
    ("✅", "LLM Reasoning Works", MINT, MINT_LT,
     "When correct code IS in candidates:",
     "81% Morph / 78% Topo accuracy"),
    ("⚠️", "Retrieval = Bottleneck", GOLD, GOLD_LT,
     "Correct code in Top-K candidates:",
     "Only 23% Morph / 19% Topo"),
    ("📊", "NOS Code Bias", ROSE, ROSE_LT,
     '"NOS" codes over-represented:',
     "8000/9 predicted in ~50% of cases"),
    ("🔍", "Strict Match Underestimates", CYAN, CYAN_LT,
     "C61 vs C61.9 fails strict match:",
     "Relaxed Topo accuracy: ~21%"),
]

positions = [(3.7, 0.3), (8.55, 0.3), (3.7, 2.15), (8.55, 2.15)]
for i, (icon, title, clr, bg, line1, line2) in enumerate(findings):
    lx, ty = positions[i]
    bx(s, lx, ty, 4.55, 1.55, fill=bg, border=clr, bw=Pt(1.5))
    # Color circle with icon
    ic = circ(s, lx + 0.2, ty + 0.2, 0.45, fill=clr)
    ic.text_frame.paragraphs[0].text = icon
    ic.text_frame.paragraphs[0].font.size = Pt(16)
    ic.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    ic.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    tx(s, lx + 0.8, ty + 0.22, 3.5, 0.3, title, sz=14, bold=True, c=clr)
    tx(s, lx + 0.8, ty + 0.6, 3.5, 0.25, line1, sz=12, c=NAVY)
    tx(s, lx + 0.8, ty + 0.9, 3.5, 0.3, line2, sz=13, bold=True, c=clr)

# Conclusion box
bx(s, 3.7, 4.0, 9.3, 1.85, fill=WHITE, border=NAVY, bw=Pt(2))
tx(s, 3.9, 4.1, 8.9, 0.3,
   "What This Means for Your Sprint", sz=16, bold=True, c=NAVY)
ml(s, 3.9, 4.5, 8.9, 1.2, [
    "→  The pipeline works — the LLM makes good coding decisions when given correct candidates.",
    "→  What's missing: a user-facing tool for clinicians to upload, assign types, trigger, review.",
    "",
    "→  That's exactly what you'll build!",
    "",
    "You don't need to fix the retrieval — just build the web tool around the existing API.",
], sz=13, c=SLATE_500)

snum(s, 5)


# ═══════════════════════════════════════════════════════════
#  SLIDE 6 – WHAT EXISTS (Sidebar + 3 status columns)
# ═══════════════════════════════════════════════════════════
s = blank()
sidebar(s, "Current\nPipeline\nComponents",
        "Research pipeline runs\nas Python scripts —\nno web interface yet", "05")

components = [
    ("OCR &\nSummarization", MINT, "✅ Complete", [
        "PaddleOCR → 77 text files",
        "Merged to 50 patient docs",
        "GPT-OSS-120B summarization",
        "Map-reduce chunking",
        "Robust multi-stage fallback",
        "Output: JSON + CSV"]),
    ("Embeddings &\nRetrieval", MINT, "✅ Complete", [
        "BGE-M3 embeddings (1024-dim)",
        "50 summary embeddings",
        "2,903 morphology embeddings",
        "1,338 topography embeddings",
        "Cosine similarity Top-K",
        "Hybrid semantic + lexical"]),
    ("RAG ICD-O\nPrediction", GOLD, "🔄 In Progress", [
        "LLM reads summary + candidates",
        "Picks best Morph + Topo code",
        "Structured JSON output",
        "3 versions (v1→v3)",
        "v3: organ-based filtering",
        "FastAPI wrapper available"]),
]

cx = 3.7
for title, clr, status, items in components:
    bx(s, cx, 0.3, 2.95, 4.8, fill=WHITE, border=WARM_200, bw=Pt(1))
    # Colored header
    bx(s, cx, 0.3, 2.95, 0.65, fill=clr)
    tx(s, cx + 0.1, 0.32, 2.75, 0.6, title.replace("\n", " "),
       sz=13, bold=True, c=WHITE, al=PP_ALIGN.CENTER)
    # Status pill
    pill_bg = MINT_LT if "Complete" in status else GOLD_LT
    pill_c = MINT if "Complete" in status else GOLD_DARK
    bx(s, cx + 0.35, 1.1, 2.25, 0.3, fill=pill_bg, border=pill_c, bw=Pt(1))
    tx(s, cx + 0.35, 1.1, 2.25, 0.3, status,
       sz=10, bold=True, c=pill_c, al=PP_ALIGN.CENTER)
    # Items
    ml(s, cx + 0.15, 1.55, 2.65, 3.4,
       [f"→  {it}" for it in items], sz=11, c=SLATE_500, sp=6)
    cx += 3.2

# Placeholders
bx(s, 3.7, 5.35, 4.5, 0.5, fill=WARM_100, border=WARM_200, bw=Pt(1))
tx(s, 3.7, 5.4, 4.5, 0.4,
   "📸 Patient Summary Example", sz=10, c=SLATE_300, al=PP_ALIGN.CENTER)
bx(s, 8.4, 5.35, 4.7, 0.5, fill=WARM_100, border=WARM_200, bw=Pt(1))
tx(s, 8.4, 5.4, 4.7, 0.4,
   "📸 ICD-O Prediction JSON Example", sz=10, c=SLATE_300, al=PP_ALIGN.CENTER)

# CTA
bx(s, 3.7, 6.05, 9.3, 0.5, fill=NAVY)
tx(s, 3.7, 6.1, 9.3, 0.4,
   "🚀  Pipeline runs as scripts — now let's make it a usable web tool!",
   sz=13, bold=True, c=GOLD, al=PP_ALIGN.CENTER)

snum(s, 6)


# ═══════════════════════════════════════════════════════════
#  SLIDE 7 – YOUR MISSION (Sidebar + 2x3 grid + scope)
# ═══════════════════════════════════════════════════════════
s = blank()
sidebar(s, "What You\nWill Build",
        "Six core features\nfor the web tool", "06")

features = [
    (1, "📤", "Document Upload", "Drag & drop medical\ndocs with preview", GOLD),
    (2, "🏷", "Type Assignment", "Assign each doc:\nArztbrief, Patho, etc.", ROSE),
    (3, "⚡", "Pipeline Trigger", "Start LLM via FastAPI\nwith type-specific prompts", CYAN),
    (4, "📊", "Results Display", "Show extracted info:\nICD-O Topo, Morph, Grade", GOLD),
    (5, "✏️", "Review & Edit", "Clinicians review,\ncorrect & confirm", ROSE),
    (6, "💾", "Export Data", "Download JSON/CSV\nfor registry integration", CYAN),
]

grid = [(3.7, 0.3), (6.8, 0.3), (9.9, 0.3),
        (3.7, 1.8), (6.8, 1.8), (9.9, 1.8)]

for i, (num, icon, title, desc, clr) in enumerate(features):
    lx, ty = grid[i]
    bx(s, lx, ty, 2.8, 1.2, fill=WHITE, border=WARM_200, bw=Pt(1))
    bx(s, lx, ty, 2.8, 0.06, fill=clr)
    nc = circ(s, lx + 0.1, ty + 0.15, 0.35, fill=clr)
    nc.text_frame.paragraphs[0].text = str(num)
    nc.text_frame.paragraphs[0].font.size = Pt(12)
    nc.text_frame.paragraphs[0].font.bold = True
    nc.text_frame.paragraphs[0].font.color.rgb = WHITE
    nc.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    nc.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    tx(s, lx + 0.55, ty + 0.12, 0.35, 0.3, icon, sz=16)
    tx(s, lx + 0.95, ty + 0.15, 1.7, 0.25, title,
       sz=12, bold=True, c=NAVY)
    ml(s, lx + 0.55, ty + 0.5, 2.1, 0.65, desc.split("\n"),
       sz=10, c=SLATE_500)

# Scope section: three columns
# YOU BUILD
bx(s, 3.7, 3.3, 3.8, 2.6, fill=WHITE, border=GOLD, bw=Pt(2))
bx(s, 3.7, 3.3, 3.8, 0.4, fill=GOLD)
tx(s, 3.7, 3.33, 3.8, 0.35, "🧑‍💻  YOU BUILD THIS",
   sz=11, bold=True, c=WHITE, al=PP_ALIGN.CENTER)
ml(s, 3.85, 3.8, 3.5, 2.0, [
    "→  Document upload & preview UI",
    "→  Document type selector / tagging",
    "→  API client for pipeline trigger",
    "→  Results display & visualization",
    "→  Review / edit workflow",
    "→  Export functionality (JSON / CSV)",
], sz=11, c=SLATE_500, sp=5)

# PROVIDED
bx(s, 7.75, 3.3, 2.8, 2.6, fill=WHITE, border=MINT, bw=Pt(2))
bx(s, 7.75, 3.3, 2.8, 0.4, fill=MINT)
tx(s, 7.75, 3.33, 2.8, 0.35, "🎁  PROVIDED",
   sz=11, bold=True, c=WHITE, al=PP_ALIGN.CENTER)
ml(s, 7.9, 3.8, 2.5, 2.0, [
    "→  FastAPI backend",
    "→  LLM summarization",
    "→  RAG retrieval",
    "→  ICD-O-3 embeddings",
    "→  Doc-type prompts",
], sz=11, c=SLATE_500, sp=5)

# OUT OF SCOPE
bx(s, 10.8, 3.3, 2.3, 2.6, fill=ROSE_LT, border=ROSE, bw=Pt(2))
bx(s, 10.8, 3.3, 2.3, 0.4, fill=ROSE)
tx(s, 10.8, 3.33, 2.3, 0.35, "❌  OUT OF SCOPE",
   sz=10, bold=True, c=WHITE, al=PP_ALIGN.CENTER)
ml(s, 10.9, 3.8, 2.1, 2.0, [
    "→  No NLP/OCR dev",
    "→  No model training",
    "→  No medical",
    "    validation",
], sz=11, c=SLATE_500, sp=5)

snum(s, 7)


# ═══════════════════════════════════════════════════════════
#  SLIDE 8 – ARCHITECTURE (Sidebar + layered diagram)
# ═══════════════════════════════════════════════════════════
s = blank()
sidebar(s, "System\nArchitecture",
        "Layered design:\nUser → Frontend →\nAPI → Backend → Data", "07")

# USER LAYER
bx(s, 3.7, 0.3, 9.3, 0.5, fill=WARM_100, border=WARM_200, bw=Pt(1))
tx(s, 3.85, 0.33, 8.9, 0.4,
   "👤  CLINICIAN / RESEARCHER — uploads documents, assigns types, reviews results",
   sz=11, bold=True, c=NAVY)

# FRONTEND LAYER
bx(s, 3.7, 1.0, 9.3, 1.8, fill=WHITE, border=GOLD, bw=Pt(2.5))
tx(s, 3.85, 1.05, 3, 0.3, "🧑‍💻  FRONTEND — You Build",
   sz=10, bold=True, c=GOLD_DARK)

fe = [("📤", "Upload"), ("🏷", "Type Select"),
      ("▶️", "Pipeline"), ("📊", "Results"),
      ("✏️", "Review"), ("💾", "Export")]
fx = 3.85
for icon, label in fe:
    bx(s, fx, 1.4, 1.4, 1.15, fill=GOLD_LT)
    tx(s, fx, 1.45, 1.4, 0.4, icon, sz=20, al=PP_ALIGN.CENTER)
    tx(s, fx, 1.85, 1.4, 0.25, label, sz=10, bold=True,
       c=GOLD_DARK, al=PP_ALIGN.CENTER)
    tx(s, fx, 2.1, 1.4, 0.25, "React/Vue/Streamlit",
       sz=8, c=SLATE_400, al=PP_ALIGN.CENTER)
    fx += 1.52

# API LAYER
bx(s, 5.5, 2.95, 5.5, 0.42, fill=NAVY, border=GOLD, bw=Pt(1))
tx(s, 5.5, 2.97, 5.5, 0.38, "⇅  FastAPI  ·  REST API  ·  JSON  ⇅",
   sz=12, bold=True, c=GOLD, al=PP_ALIGN.CENTER)

# BACKEND LAYER
bx(s, 3.7, 3.6, 9.3, 1.6, fill=WHITE, border=MINT, bw=Pt(2.5))
tx(s, 3.85, 3.65, 3, 0.3, "🎁  BACKEND — Provided",
   sz=10, bold=True, c=MINT)

be = [("📄", "Doc Processing", "OCR + merge"),
      ("🧠", "LLM Summary", "GPT-OSS-120B"),
      ("📐", "RAG Retrieval", "BGE-M3 Top-K"),
      ("🏷", "ICD-O Assign", "Morph + Topo")]
bex = 4.0
for icon, title, desc in be:
    bx(s, bex, 4.0, 2.0, 0.95, fill=MINT_LT)
    tx(s, bex, 4.05, 2.0, 0.35, f"{icon}  {title}",
       sz=11, bold=True, c=MINT, al=PP_ALIGN.CENTER)
    tx(s, bex, 4.4, 2.0, 0.25, desc, sz=10, c=SLATE_500, al=PP_ALIGN.CENTER)
    bex += 2.15

# DATA LAYER
bx(s, 3.7, 5.45, 9.3, 0.45, fill=NAVY)
tx(s, 3.7, 5.48, 9.3, 0.38,
   "📦  ICD-O-3 Dictionary (4,241 terms)  ·  BGE-M3 Embeddings  ·  "
   "Patient OCR Docs (50 patients)",
   sz=10, c=SLATE_400, al=PP_ALIGN.CENTER)

snum(s, 8)


# ═══════════════════════════════════════════════════════════
#  SLIDE 9 – TIMELINE (Sidebar + horizontal day blocks)
# ═══════════════════════════════════════════════════════════
s = blank()
sidebar(s, "Sprint Plan\n& Team\nRoles",
        "4 working days\nfrom kick-off to demo", "08")

# Horizontal day blocks (different from dot-timeline and Gantt!)
days = [
    ("FRI 20", "Afternoon", GOLD, [
        "Kick-off & team formation",
        "Understand pipeline & API",
        "Set up dev environment",
    ]),
    ("MON 23", "Full Day", ROSE, [
        "Build upload module",
        "Document type selector",
        "Connect to FastAPI backend",
    ]),
    ("TUE 24", "Full Day", CYAN, [
        "Results display panel",
        "Review / edit interface",
        "Export & progress tracking",
    ]),
    ("WED 25", "AM → PM 🏆", MINT, [
        "Polish UI & bug fixes",
        "Prepare demo",
        "Final presentation!",
    ]),
]

dx = 3.7
for day, time, clr, tasks in days:
    bx(s, dx, 0.3, 2.2, 3.8, fill=WHITE, border=WARM_200, bw=Pt(1))
    # Day header
    bx(s, dx, 0.3, 2.2, 0.7, fill=clr)
    tx(s, dx, 0.32, 2.2, 0.35, day, sz=16, bold=True, c=WHITE, al=PP_ALIGN.CENTER)
    tx(s, dx, 0.65, 2.2, 0.25, time, sz=10, c=WHITE, al=PP_ALIGN.CENTER)
    # Tasks
    ml(s, dx + 0.15, 1.15, 1.9, 2.8,
       [f"→  {t}" for t in tasks], sz=11, c=SLATE_500, sp=8)
    dx += 2.4

# Team roles section
tx(s, 3.7, 4.4, 5, 0.3, "Team Roles", sz=16, bold=True, c=NAVY)

roles = [
    ("🎨", "Frontend Dev", "UI components & styling", GOLD),
    ("⚙️", "API Lead", "FastAPI integration & data flow", ROSE),
    ("🩺", "Clinical Expert", "ICD-O understanding & UX", CYAN),
    ("🏆", "Demo Lead", "Presentation & live demo", MINT),
]

rx = 3.7
for icon, title, desc, clr in roles:
    bx(s, rx, 4.85, 2.2, 0.9, fill=WHITE, border=clr, bw=Pt(1.5))
    ic = circ(s, rx + 0.15, 4.95, 0.4, fill=clr)
    ic.text_frame.paragraphs[0].text = icon
    ic.text_frame.paragraphs[0].font.size = Pt(14)
    ic.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    ic.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    tx(s, rx + 0.65, 4.95, 1.4, 0.25, title,
       sz=12, bold=True, c=NAVY)
    tx(s, rx + 0.65, 5.25, 1.4, 0.25, desc, sz=9, c=SLATE_500)
    rx += 2.4

snum(s, 9)


# ═══════════════════════════════════════════════════════════
#  SLIDE 10 – CTA (Full dark closing — matches title)
# ═══════════════════════════════════════════════════════════
s = blank()

# Full navy background
bx(s, 0, 0, 13.333, 7.5, fill=NAVY)
bx(s, 0, 0, 13.333, 0.06, fill=GOLD)

# Decorative
circ(s, -1.0, -1.0, 3.5, fill=NAVY_MID)
circ(s, 11.5, 5.5, 3.0, fill=NAVY_MID)

tx(s, 1.0, 0.5, 11.3, 0.5, "What You'll Gain",
   sz=32, bold=True, c=WHITE, al=PP_ALIGN.CENTER)

# Three skill cards
skills = [
    ("🏥", "Clinical Knowledge", GOLD, [
        "ICD-O oncology coding",
        "Clinical document types",
        "Tumor documentation workflow",
        "Structured vs unstructured data"]),
    ("💻", "Technical Skills", ROSE, [
        "FastAPI integration",
        "Modern frontend dev",
        "LLM pipeline architecture",
        "RAG system design"]),
    ("🎯", "Professional Growth", CYAN, [
        "Agile sprint workflow",
        "Team collaboration",
        "Product demo & pitching",
        "Digital Health project"]),
]

sx = 1.5
for icon, title, clr, items in skills:
    bx(s, sx, 1.3, 3.2, 2.2, fill=NAVY_MID, border=clr, bw=Pt(1.5))
    bx(s, sx, 1.3, 3.2, 0.06, fill=clr)
    tx(s, sx + 0.15, 1.4, 0.4, 0.35, icon, sz=20, c=WHITE)
    tx(s, sx + 0.6, 1.42, 2.4, 0.3, title, sz=14, bold=True, c=clr)
    ml(s, sx + 0.15, 1.85, 2.9, 1.5,
       [f"→  {it}" for it in items], sz=11, c=SLATE_400, sp=5)
    sx += 3.5

# CTA
tx(s, 1.0, 3.9, 11.3, 0.5,
   "Ready to Transform Oncology Documentation?",
   sz=26, bold=True, c=GOLD, al=PP_ALIGN.CENTER)

# Feature pills
pills = [
    ("🧬  Real clinical data", GOLD),
    ("🤖  LLM + RAG pipeline", ROSE),
    ("🏥  Impact on registries", CYAN),
    ("🚀  Research → Production", MINT),
]
px_p = 1.8
for text, clr in pills:
    bx(s, px_p, 4.6, 2.3, 0.38, fill=NAVY_MID, border=clr, bw=Pt(1))
    tx(s, px_p, 4.62, 2.3, 0.35, text,
       sz=10, bold=True, c=clr, al=PP_ALIGN.CENTER)
    px_p += 2.5

# Supervisors
bx(s, 3.5, 5.3, 6.3, 0.9, fill=NAVY_MID, border=NAVY_LIGHT, bw=Pt(1))
iw2 = circ(s, 4.3, 5.45, 0.5, fill=GOLD)
iw2.text_frame.paragraphs[0].text = "IW"
iw2.text_frame.paragraphs[0].font.size = Pt(13)
iw2.text_frame.paragraphs[0].font.bold = True
iw2.text_frame.paragraphs[0].font.color.rgb = NAVY
iw2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
iw2.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

nd2 = circ(s, 4.9, 5.45, 0.5, fill=ROSE)
nd2.text_frame.paragraphs[0].text = "ND"
nd2.text_frame.paragraphs[0].font.size = Pt(13)
nd2.text_frame.paragraphs[0].font.bold = True
nd2.text_frame.paragraphs[0].font.color.rgb = WHITE
nd2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
nd2.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

tx(s, 5.6, 5.4, 3.8, 0.3,
   "Isabella Catharina Wiest & Naghme Dashti", sz=13, bold=True, c=WHITE)
tx(s, 5.6, 5.73, 3.8, 0.25,
   "EKFZ Digital Health · TU Dresden", sz=10, c=SLATE_400)

# Logo placeholders
for i, lb in enumerate(["EKFZ", "TUD", "CGC"]):
    lx = 10.5 + i * 1.0
    bx(s, lx, 6.5, 0.8, 0.55, fill=NAVY_MID, border=NAVY_LIGHT, bw=Pt(1))
    tx(s, lx, 6.55, 0.8, 0.45, lb, sz=8, c=SLATE_400, al=PP_ALIGN.CENTER)

snum(s, 10)


# ═════════════════════════════════════════════════════════════
output = "Tumor_PitchDeck_DesignC.pptx"
prs.save(output)
print(f"\n✅ Saved: {output}")
print(f"   Slides: {len(prs.slides)}")
print(f"\n🎨 Design C — 'Dark Sidebar + Warm Gold':")
print(f"   • Dark navy sidebar on every content slide")
print(f"   • Gold accent line separator")
print(f"   • Warm palette: Navy + Gold + Rose + Cyan + Mint")
print(f"   • Full dark title + closing slides")
print(f"   • Top-bar colored cards (not left-accent)")
print(f"   • Vertical waterfall pipeline (not horizontal)")
print(f"   • Horizontal day-block timeline (not dots/Gantt)")
print(f"   • Layered architecture diagram")
