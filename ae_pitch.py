# ============================================================
#  AE Coding Pitch Deck – Clinicum Digitale Spring 2026
#  Run: pip install python-pptx   then   python ae_pitch.py
# ============================================================

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colours ──────────────────────────────────────────────────
TEAL       = RGBColor(0x0D, 0x94, 0x88)
TEAL_DARK  = RGBColor(0x0F, 0x76, 0x6E)
TEAL_LIGHT = RGBColor(0xCC, 0xFB, 0xF1)
PURPLE     = RGBColor(0x63, 0x66, 0xF1)
PURPLE_LT  = RGBColor(0xEE, 0xF2, 0xFF)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_LT   = RGBColor(0xFF, 0xFB, 0xEB)
RED        = RGBColor(0xEF, 0x44, 0x44)
RED_LT     = RGBColor(0xFE, 0xF2, 0xF2)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG_SUBTLE  = RGBColor(0xF1, 0xF5, 0xF9)
BG_LIGHT   = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_PRI   = RGBColor(0x0F, 0x17, 0x2A)
TEXT_SEC   = RGBColor(0x47, 0x55, 0x69)
TEXT_DIM   = RGBColor(0x94, 0xA3, 0xB8)
BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
BLUE       = RGBColor(0x25, 0x63, 0xEB)
BLUE_LT    = RGBColor(0xEF, 0xF6, 0xFF)

# ── Presentation Setup ───────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# ── Helper functions ─────────────────────────────────────────
def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text,
                font_size=14, bold=False, color=TEXT_PRI,
                alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_textbox(slide, left, top, width, height, runs,
                     alignment=PP_ALIGN.LEFT, line_spacing=1.3):
    """runs = list of dicts: {text, size, bold, color, font_name}"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing * 14)
    for i, r in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
        else:
            run = p.add_run()
        run.text = r.get("text", "")
        run.font.size = Pt(r.get("size", 14))
        run.font.bold = r.get("bold", False)
        run.font.color.rgb = r.get("color", TEXT_PRI)
        run.font.name = r.get("font_name", "Calibri")
    return txBox

def add_rect(slide, left, top, width, height, fill_color=BG_SUBTLE,
             border_color=None, border_width=Pt(1), corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=13, color=TEXT_SEC, bold=False,
                          alignment=PP_ALIGN.LEFT, line_spacing=1.4,
                          font_name="Calibri"):
    """lines = list of strings, each becomes a paragraph"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
    return txBox

def add_badge(slide, left, top, text, bg_color=TEAL_LIGHT, text_color=TEAL_DARK):
    w, h = 3.5, 0.4
    shape = add_rect(slide, left, top, w, h, fill_color=bg_color)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(10)
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.color.rgb = text_color
    shape.text_frame.paragraphs[0].font.name = "Calibri"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_footer(slide, text="EKFZ Digital Health  ·  TU Dresden  ·  March 2026"):
    add_textbox(slide, 0.6, 7.0, 5, 0.3, text,
                font_size=9, color=TEXT_DIM)

def add_card_box(slide, left, top, width, height, title, body_lines,
                 icon="", border_color=BORDER):
    add_rect(slide, left, top, width, height,
             fill_color=WHITE, border_color=border_color, border_width=Pt(1.5))
    y = top + 0.2
    if icon:
        add_textbox(slide, left + 0.25, y, 0.5, 0.4, icon, font_size=22)
        y += 0.35
    add_textbox(slide, left + 0.25, y, width - 0.5, 0.35, title,
                font_size=15, bold=True, color=TEXT_PRI)
    y += 0.4
    add_multiline_textbox(slide, left + 0.25, y, width - 0.5,
                          height - (y - top) - 0.2, body_lines,
                          font_size=12, color=TEXT_SEC)

def add_pipeline_box(slide, left, top, w, h, num, title, desc, active=True):
    bc = TEAL if active else BORDER
    add_rect(slide, left, top, w, h, fill_color=WHITE,
             border_color=bc, border_width=Pt(2))
    # number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left + w/2 - 0.2), Inches(top + 0.2),
        Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = TEAL
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = str(num)
    circle.text_frame.paragraphs[0].font.size = Pt(12)
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(slide, left + 0.15, top + 0.7, w - 0.3, 0.35, title,
                font_size=14, bold=True, color=TEXT_PRI,
                alignment=PP_ALIGN.CENTER)
    add_multiline_textbox(slide, left + 0.15, top + 1.05, w - 0.3, h - 1.25,
                          desc, font_size=11, color=TEXT_SEC,
                          alignment=PP_ALIGN.CENTER)

def add_arrow_text(slide, left, top, text="→"):
    add_textbox(slide, left, top, 0.4, 0.5, text,
                font_size=24, bold=True, color=TEAL,
                alignment=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════
#  SLIDE 1 – TITLE
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()

# subtle decorative circle (top-right)
circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1),
                            Inches(4), Inches(4))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0xE6, 0xFB, 0xF7)
circle.line.fill.background()
circle.fill.fore_color.brightness = 0.1

add_badge(s, 0.7, 0.7,
          "CLINICUM DIGITALE  ·  SPRING 2026  ·  PROJECT X",
          TEAL_LIGHT, TEAL_DARK)

add_rich_textbox(s, 0.7, 1.4, 9, 1.5, [
    {"text": "LLM-Powered ", "size": 42, "bold": True, "color": TEXT_PRI},
    {"text": "MedDRA\n", "size": 42, "bold": True, "color": TEAL},
    {"text": "Adverse Event Coding", "size": 42, "bold": True, "color": TEXT_PRI},
])

add_textbox(s, 0.7, 3.2, 8, 0.8,
    "Build a production-ready web tool that transforms free-text adverse "
    "events into standardized MedDRA codes — powered by Retrieval-Augmented Reasoning.",
    font_size=17, color=TEXT_SEC)

# Supervisor box
add_rect(s, 0.7, 4.5, 7, 0.9, fill_color=BG_SUBTLE)
add_textbox(s, 1.0, 4.6, 6.5, 0.3,
            "Isabella Catharina Wiest & Naghme Dashti",
            font_size=16, bold=True, color=TEXT_PRI)
add_textbox(s, 1.0, 4.95, 6.5, 0.3,
            "Else Kröner Fresenius Center for Digital Health  ·  TU Dresden",
            font_size=12, color=TEXT_DIM)

# Logo placeholders
add_rect(s, 9.5, 6.2, 1.2, 0.8, fill_color=BG_SUBTLE, border_color=BORDER)
add_textbox(s, 9.5, 6.3, 1.2, 0.6, "EKFZ\nLogo", font_size=9,
            color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_rect(s, 10.9, 6.2, 1.2, 0.8, fill_color=BG_SUBTLE, border_color=BORDER)
add_textbox(s, 10.9, 6.3, 1.2, 0.6, "TUD\nLogo", font_size=9,
            color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_rect(s, 12.3, 6.2, 0.8, 0.8, fill_color=BG_SUBTLE, border_color=BORDER)
add_textbox(s, 12.3, 6.3, 0.8, 0.6, "CGC\nLogo", font_size=8,
            color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 2 – THE PROBLEM
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
add_badge(s, 0.7, 0.5, "THE CHALLENGE", PURPLE_LT, PURPLE)
add_textbox(s, 0.7, 1.0, 10, 0.6, "Why Does This Matter?",
            font_size=32, bold=True, color=TEXT_PRI)

# Problem points (left side)
problems = [
    ("💊  Adverse Events = backbone of drug safety",
     "Every clinical trial generates hundreds to thousands of AE reports requiring coding."),
    ("📝  AEs are written as free text",
     'Heterogeneous, ambiguous descriptions: "blurry vision", "difficulty breathing", "Bilirubin Erhoehung"'),
    ("⏱️  Manual MedDRA coding is a bottleneck",
     "Time-consuming, subjective, prone to inter-rater variability, hard to scale."),
    ("🔍  80,000+ possible codes in MedDRA",
     "Finding the right LLT is like a needle in a haystack."),
]

y = 1.75
for title, desc in problems:
    add_rect(s, 0.7, y, 6.0, 0.72, fill_color=WHITE, border_color=BORDER)
    add_textbox(s, 0.9, y + 0.05, 5.6, 0.3, title,
                font_size=13, bold=True, color=TEXT_PRI)
    add_textbox(s, 0.9, y + 0.35, 5.6, 0.35, desc,
                font_size=11, color=TEXT_SEC)
    y += 0.85

# Current vs Vision (right side)
# Current workflow
add_rect(s, 7.2, 1.75, 5.4, 1.2, fill_color=RED_LT, border_color=RED,
         border_width=Pt(1.5))
add_textbox(s, 7.4, 1.85, 5, 0.3, "❌  Current Workflow",
            font_size=14, bold=True, color=RED)
add_textbox(s, 7.4, 2.2, 5, 0.3,
            "📝 Free-text AE  →  👨‍⚕️ Manual Coding  →  📋 MedDRA Code",
            font_size=12, color=TEXT_SEC)
add_textbox(s, 7.4, 2.55, 5, 0.3,
            "⏳ Hours per trial  ·  ❌ Inconsistent  ·  📉 Doesn't scale",
            font_size=11, bold=True, color=RED)

# Our vision
add_rect(s, 7.2, 3.15, 5.4, 1.2, fill_color=TEAL_LIGHT, border_color=TEAL,
         border_width=Pt(1.5))
add_textbox(s, 7.4, 3.25, 5, 0.3, "✅  Our Vision",
            font_size=14, bold=True, color=TEAL_DARK)
add_textbox(s, 7.4, 3.6, 5, 0.3,
            "📝 Free-text AE  →  🤖 LLM + RAG  →  ✓ Human Review",
            font_size=12, color=TEXT_SEC)
add_textbox(s, 7.4, 3.95, 5, 0.3,
            "⚡ Seconds per AE  ·  🎯 91–97% correct  ·  📈 Scalable",
            font_size=11, bold=True, color=TEAL_DARK)

# Quote box
add_rect(s, 7.2, 4.6, 5.4, 1.3, fill_color=BG_SUBTLE)
add_multiline_textbox(s, 7.4, 4.7, 5.0, 1.1, [
    "\"The current manual coding process is not",
    "just slow — it introduces systematic bias",
    "and reduces reproducibility across trials.\"",
], font_size=12, color=TEXT_SEC)

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 3 – MedDRA HIERARCHY
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
add_badge(s, 0.7, 0.5, "BACKGROUND", TEAL_LIGHT, TEAL_DARK)
add_textbox(s, 0.7, 1.0, 10, 0.6, "The MedDRA Hierarchy",
            font_size=32, bold=True, color=TEXT_PRI)
add_textbox(s, 0.7, 1.5, 10, 0.35,
            "Medical Dictionary for Regulatory Activities — 5-level hierarchical terminology by ICH",
            font_size=14, color=TEXT_SEC)

# Hierarchy boxes (left)
levels = [
    ("SOC — System Organ Class", "27 categories · e.g., \"Nervous system disorders\"",
     PURPLE_LT, PURPLE, 4.8),
    ("PT — Preferred Term", "~22,000 terms · e.g., \"Headache\"",
     BLUE_LT, BLUE, 4.2),
    ("LLT — Low-Level Term", "~80,000 terms · e.g., \"Throbbing headache\"",
     TEAL_LIGHT, TEAL_DARK, 3.6),
]

y = 2.2
for title, desc, bg, tc, w in levels:
    cx = 0.7 + (4.8 - w) / 2
    add_rect(s, cx, y, w, 0.75, fill_color=bg, border_color=tc, border_width=Pt(1.5))
    add_textbox(s, cx + 0.15, y + 0.08, w - 0.3, 0.3, title,
                font_size=12, bold=True, color=tc, alignment=PP_ALIGN.CENTER)
    add_textbox(s, cx + 0.15, y + 0.38, w - 0.3, 0.3, desc,
                font_size=11, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)
    y += 0.85
    if y < 4.6:
        add_textbox(s, 2.7, y - 0.12, 0.6, 0.25, "▼",
                    font_size=16, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# Examples table (right)
add_textbox(s, 6.5, 2.1, 6, 0.35, "Real-World Examples",
            font_size=16, bold=True, color=TEAL_DARK)

# Table header
add_rect(s, 6.5, 2.55, 6.2, 0.45, fill_color=BG_SUBTLE)
headers = [("Free-text AE", 6.5, 1.6), ("LLT", 8.1, 1.4),
           ("PT", 9.5, 1.3), ("SOC", 10.8, 1.9)]
for text, x, w in headers:
    add_textbox(s, x, 2.58, w, 0.35, text,
                font_size=10, bold=True, color=TEAL_DARK)

# Table rows
rows = [
    ('"blurry vision"', "Blurred vision", "Vision blurred", "Eye disorders"),
    ('"difficulty breathing"', "Shortness of breath", "Dyspnea", "Respiratory disorders"),
    ('"Skin rash after injection"', "Injection site rash", "Inj. site reaction", "Skin disorders"),
]
y = 3.1
for ae, llt, pt, soc in rows:
    if rows.index((ae, llt, pt, soc)) % 2 == 0:
        add_rect(s, 6.5, y - 0.05, 6.2, 0.5, fill_color=WHITE)
    add_textbox(s, 6.5, y, 1.6, 0.4, ae, font_size=11, bold=True, color=TEXT_PRI)
    add_textbox(s, 8.1, y, 1.4, 0.4, llt, font_size=11, color=TEXT_SEC)
    add_textbox(s, 9.5, y, 1.3, 0.4, pt, font_size=11, color=TEXT_SEC)
    add_textbox(s, 10.8, y, 1.9, 0.4, soc, font_size=11, color=TEXT_SEC)
    y += 0.5

# Key challenge box
add_rect(s, 6.5, 4.8, 6.2, 0.8, fill_color=TEAL_LIGHT)
add_rect(s, 6.5, 4.8, 0.06, 0.8, fill_color=TEAL)
add_multiline_textbox(s, 6.75, 4.85, 5.8, 0.7, [
    "Key Challenge: Multiple LLTs can be clinically",
    'acceptable for the same AE — e.g., "Headache"',
    'vs "Headache NOS" vs "Cephalalgia"',
], font_size=12, color=TEXT_PRI)

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 4 – RAG PIPELINE
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
add_badge(s, 0.7, 0.5, "OUR APPROACH", TEAL_LIGHT, TEAL_DARK)
add_textbox(s, 0.7, 1.0, 10, 0.6, "Retrieval-Augmented Reasoning Pipeline",
            font_size=32, bold=True, color=TEXT_PRI)
add_textbox(s, 0.7, 1.55, 10, 0.3,
            "Three-stage architecture: semantic retrieval → structured LLM reasoning → deterministic mapping",
            font_size=14, color=TEXT_SEC)

# Three pipeline stages
stage_data = [
    (1, "Candidate Generation",
     ["Dense semantic retrieval", "MiniLM embeddings (384-dim)",
      "Top-100 similar LLTs", "via cosine similarity"]),
    (2, "LLM Reasoning",
     ["Chain-of-Thought prompting", "LLaMA-3.3-70B-Instruct",
      "Structured comparison", 'Output: "Final answer: <LLT>"']),
    (3, "Hierarchical Mapping",
     ["Deterministic lookup", "LLT → PT → SOC",
      "MedDRA v25.0", "Traceable & auditable"]),
]

x = 1.0
for num, title, desc in stage_data:
    add_pipeline_box(s, x, 2.3, 3.3, 2.6, num, title, desc)
    if num < 3:
        add_arrow_text(s, x + 3.35, 3.4)
    x += 3.8

# Bottom highlight
add_rect(s, 0.7, 5.3, 12, 0.8, fill_color=TEAL_LIGHT)
add_rect(s, 0.7, 5.3, 0.06, 0.8, fill_color=TEAL)
add_textbox(s, 1.0, 5.4, 11.5, 0.6,
    "💡  The LLM pipeline already exists as a FastAPI backend — "
    "your mission is to build the web application & user workflow around it!",
    font_size=15, bold=True, color=TEAL_DARK, alignment=PP_ALIGN.CENTER)

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 5 – PERFORMANCE RESULTS
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
add_badge(s, 0.7, 0.5, "VALIDATED RESULTS", PURPLE_LT, PURPLE)
add_textbox(s, 0.7, 1.0, 10, 0.6,
            "Performance Across 3 Clinical Trials",
            font_size=32, bold=True, color=TEXT_PRI)
add_textbox(s, 0.7, 1.55, 10, 0.3,
            "Tested on AML oncology datasets — validated by blinded expert clinical review",
            font_size=14, color=TEXT_SEC)

# Stat boxes
stats = [
    ("91–97%", "Clinical Correctness"),
    ("78–85%", "PT-Level Accuracy"),
    ("90–93%", "SOC-Level Accuracy"),
    ("50–58%", "Exact LLT Match"),
]
x = 0.7
for value, label in stats:
    add_rect(s, x, 2.1, 2.8, 1.1, fill_color=WHITE, border_color=BORDER,
             border_width=Pt(1.5))
    add_textbox(s, x, 2.2, 2.8, 0.5, value,
                font_size=30, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x, 2.75, 2.8, 0.35, label,
                font_size=11, bold=True, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    x += 3.05

# RAG vs Zero-shot comparison (left)
add_textbox(s, 0.7, 3.5, 5, 0.35, "RAG Pipeline vs Zero-Shot",
            font_size=16, bold=True, color=TEXT_PRI)

comparisons = [
    ("PT Accuracy", [("RAG", 85, TEAL), ("Zero-shot", 72, PURPLE)]),
    ("SOC Accuracy", [("RAG", 92, TEAL), ("Zero-shot", 74, PURPLE)]),
]
y = 3.95
for section, bars in comparisons:
    add_textbox(s, 0.7, y, 5, 0.25, section,
                font_size=11, bold=True, color=TEXT_DIM)
    y += 0.3
    for label, pct, clr in bars:
        add_textbox(s, 0.7, y, 1.0, 0.3, label, font_size=11, color=TEXT_SEC)
        # bar background
        add_rect(s, 1.7, y + 0.03, 4.5, 0.3, fill_color=BG_SUBTLE)
        # bar fill
        fill_w = 4.5 * pct / 100
        add_rect(s, 1.7, y + 0.03, fill_w, 0.3, fill_color=clr)
        add_textbox(s, 1.7 + fill_w - 0.6, y + 0.03, 0.6, 0.3,
                    f"{pct}%", font_size=11, bold=True, color=WHITE,
                    alignment=PP_ALIGN.RIGHT)
        y += 0.38
    y += 0.15

# Right side: why LLT exact ≠ clinical quality
add_rect(s, 7.0, 3.5, 5.6, 3.1, fill_color=WHITE, border_color=BORDER,
         border_width=Pt(1.5))
add_textbox(s, 7.2, 3.6, 5.2, 0.35, "Why LLT Exact ≠ Clinical Quality",
            font_size=16, bold=True, color=TEXT_PRI)
add_multiline_textbox(s, 7.2, 4.05, 5.2, 2.4, [
    "📊  38–41% of 'errors' under strict LLT matching were",
    "     actually clinically acceptable alternatives",
    "",
    "🩺  Multiple LLTs can be valid for the same AE",
    "",
    "✅  Blinded expert review confirmed",
    "     real-world clinical usability",
    "",
    "→ That's why we need Human-in-the-Loop!",
], font_size=13, color=TEXT_SEC)

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 6 – EXISTING PROTOTYPES
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
add_badge(s, 0.7, 0.5, "WHAT EXISTS", TEAL_LIGHT, TEAL_DARK)
add_textbox(s, 0.7, 1.0, 10, 0.6, "Existing Prototypes",
            font_size=32, bold=True, color=TEXT_PRI)
add_textbox(s, 0.7, 1.55, 10, 0.3,
            "We've built Flask-based demos — your job is to create the production-ready tool",
            font_size=14, color=TEXT_SEC)

# Prototype 1
add_card_box(s, 0.7, 2.1, 5.8, 3.6,
    "🎨  MedDRA RAG Assistant",
    ["• Single AE text input → instant prediction",
     "• Shows LLT / PT / SOC with MedDRA codes",
     "• Top candidate terms with similarity scores",
     "• MedDRA hierarchy tree view",
     "• Latency tracking & fallback indicator",
     "• Quick example buttons for testing",
     "",
     "✅ Working prototype · Flask + FastAPI"],
    border_color=TEAL)

# Prototype 2
add_card_box(s, 6.8, 2.1, 5.8, 3.6,
    "📋  Clinical Coding Board",
    ["• Batch CSV/XLSX upload with progress bar",
     "• Row-by-row LLM prediction",
     "• Human-in-the-Loop: Accept / Reject",
     "• Top-5 alternative suggestions dropdown",
     "• Manual LLT override option",
     "• Export reviewed results as Excel",
     "",
     "✅ Working prototype · Flask + FastAPI"],
    border_color=TEAL)

# Screenshot placeholders
add_rect(s, 1.0, 5.9, 5.2, 0.5, fill_color=BG_SUBTLE, border_color=BORDER)
add_textbox(s, 1.0, 5.95, 5.2, 0.4,
            "📸 Screenshot Placeholder — Demo 1",
            font_size=10, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_rect(s, 7.1, 5.9, 5.2, 0.5, fill_color=BG_SUBTLE, border_color=BORDER)
add_textbox(s, 7.1, 5.95, 5.2, 0.4,
            "📸 Screenshot Placeholder — Demo 2",
            font_size=10, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# Bottom CTA
add_rect(s, 0.7, 6.55, 12, 0.55, fill_color=PURPLE_LT)
add_textbox(s, 0.7, 6.6, 12, 0.45,
    "🚀  These prototypes prove the concept works — now let's build it right!",
    font_size=15, bold=True, color=PURPLE, alignment=PP_ALIGN.CENTER)

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 7 – YOUR MISSION
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
add_badge(s, 0.7, 0.5, "YOUR MISSION", PURPLE_LT, PURPLE)
add_textbox(s, 0.7, 1.0, 10, 0.6, "What You Will Build 🛠️",
            font_size=32, bold=True, color=TEXT_PRI)

# Left: Feature list
features = [
    ("📤", "CSV/XLSX Upload",
     "Drag & drop adverse event files with column mapping"),
    ("⚡", "FastAPI Integration",
     "Connect to existing LLM backend (black-box) via REST API"),
    ("📊", "Results Dashboard",
     "Display LLT/PT/SOC predictions with Top-5 alternatives"),
    ("✅", "Human-in-the-Loop Validation",
     "Accept / Reject / Modify each prediction interactively"),
    ("📈", "Progress Tracking + Export",
     "Real-time progress bar, intermediate saves, CSV/Excel export"),
]
y = 1.7
for icon, title, desc in features:
    add_rect(s, 0.7, y, 6.8, 0.7, fill_color=WHITE, border_color=BORDER)
    add_textbox(s, 0.85, y + 0.07, 0.4, 0.35, icon, font_size=20)
    add_textbox(s, 1.35, y + 0.07, 3, 0.3, title,
                font_size=14, bold=True, color=TEXT_PRI)
    add_textbox(s, 1.35, y + 0.35, 5.8, 0.3, desc,
                font_size=12, color=TEXT_SEC)
    y += 0.78

# Right: Build vs Provided
# You Build
add_rect(s, 8.0, 1.7, 4.6, 2.3, fill_color=WHITE,
         border_color=TEAL, border_width=Pt(2))
add_rect(s, 8.15, 1.85, 1.7, 0.3, fill_color=TEAL_LIGHT)
add_textbox(s, 8.15, 1.85, 1.7, 0.3, "🧑‍💻 YOU BUILD",
            font_size=9, bold=True, color=TEAL_DARK, alignment=PP_ALIGN.CENTER)
add_multiline_textbox(s, 8.2, 2.25, 4.2, 1.6, [
    "• Upload module & file parsing",
    "• UI/UX for review workflow",
    "• API client for FastAPI calls",
    "• State management & progress",
    "• Export functionality",
    "• Responsive, modern design",
], font_size=12, color=TEXT_SEC)

# Provided
add_rect(s, 8.0, 4.15, 4.6, 1.6, fill_color=WHITE,
         border_color=PURPLE, border_width=Pt(2))
add_rect(s, 8.15, 4.3, 1.7, 0.3, fill_color=PURPLE_LT)
add_textbox(s, 8.15, 4.3, 1.7, 0.3, "🎁 PROVIDED",
            font_size=9, bold=True, color=PURPLE, alignment=PP_ALIGN.CENTER)
add_multiline_textbox(s, 8.2, 4.7, 4.2, 1.0, [
    "• FastAPI backend (LLM pipeline)",
    "• MiniLM embeddings + retrieval",
    "• LLaMA-3.3-70B inference",
    "• MedDRA v25.0 data resources",
], font_size=12, color=TEXT_SEC)

# Out of scope
add_rect(s, 8.0, 5.9, 4.6, 0.8, fill_color=RED_LT,
         border_color=RED, border_width=Pt(1.5))
add_textbox(s, 8.2, 5.95, 4.2, 0.3, "❌ Out of Scope",
            font_size=12, bold=True, color=RED)
add_textbox(s, 8.2, 6.3, 4.2, 0.35,
            "No model training, no NLP dev, no medical validation",
            font_size=11, color=TEXT_SEC)

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 8 – ARCHITECTURE
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
add_badge(s, 0.7, 0.5, "TECHNICAL ARCHITECTURE", TEAL_LIGHT, TEAL_DARK)
add_textbox(s, 0.7, 1.0, 10, 0.6, "System Overview",
            font_size=32, bold=True, color=TEXT_PRI)

# Frontend box (you build)
add_rect(s, 1.0, 1.8, 11.3, 1.8, fill_color=WHITE,
         border_color=TEAL, border_width=Pt(2))
add_rect(s, 1.15, 1.95, 1.8, 0.3, fill_color=TEAL_LIGHT)
add_textbox(s, 1.15, 1.95, 1.8, 0.3, "🧑‍💻 FRONTEND — You Build",
            font_size=9, bold=True, color=TEAL_DARK, alignment=PP_ALIGN.CENTER)

# Frontend modules
fe_modules = [
    ("📤", "Upload Module", "CSV/XLSX parsing"),
    ("🔍", "Review Dashboard", "Accept / Reject / Modify"),
    ("💾", "Export Module", "CSV/Excel download"),
]
x = 1.5
for icon, title, desc in fe_modules:
    add_rect(s, x, 2.4, 3.2, 0.95, fill_color=BG_SUBTLE)
    add_textbox(s, x + 0.1, 2.45, 3.0, 0.35, f"{icon}  {title}",
                font_size=14, bold=True, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + 0.1, 2.8, 3.0, 0.35, desc,
                font_size=11, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    x += 3.5

# Arrow
add_textbox(s, 5.5, 3.7, 2.5, 0.5, "⇅  REST API  ⇅",
            font_size=18, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)

# Backend box (provided)
add_rect(s, 1.0, 4.3, 11.3, 1.8, fill_color=WHITE,
         border_color=PURPLE, border_width=Pt(2))
add_rect(s, 1.15, 4.45, 2.2, 0.3, fill_color=PURPLE_LT)
add_textbox(s, 1.15, 4.45, 2.2, 0.3, "🎁 BACKEND — Provided",
            font_size=9, bold=True, color=PURPLE, alignment=PP_ALIGN.CENTER)

be_modules = [
    ("🎯", "MiniLM Embeddings", "384-dim retrieval"),
    ("🧠", "LLaMA-3.3-70B", "Chain-of-Thought"),
    ("📖", "MedDRA v25.0", "LLT → PT → SOC"),
]
x = 1.5
for icon, title, desc in be_modules:
    add_rect(s, x, 4.9, 3.2, 0.95, fill_color=BG_SUBTLE)
    add_textbox(s, x + 0.1, 4.95, 3.0, 0.35, f"{icon}  {title}",
                font_size=14, bold=True, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + 0.1, 5.3, 3.0, 0.35, desc,
                font_size=11, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    x += 3.5

# Tech stack labels
tech = ["🖥 Frontend: React / Vue / Streamlit",
        "⚡ API: FastAPI (provided)",
        "🧠 LLM: LLaMA-3.3-70B",
        "📦 Data: MedDRA CSV"]
x = 1.0
for t in tech:
    add_rect(s, x, 6.4, 2.8, 0.45, fill_color=BG_SUBTLE)
    add_textbox(s, x, 6.45, 2.8, 0.35, t,
                font_size=11, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)
    x += 3.0

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 9 – TIMELINE & TEAM
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
add_badge(s, 0.7, 0.5, "SPRINT PLAN", AMBER_LT, RGBColor(0xB4, 0x53, 0x09))
add_textbox(s, 0.7, 1.0, 10, 0.6, "Project Timeline & Team Roles",
            font_size=32, bold=True, color=TEXT_PRI)

# Timeline (left)
timeline = [
    ("Fri, 20 March — PM", "Kick-off · Understand pipeline & API docs\nSet up dev environment · Team roles"),
    ("Mon, 23 March", "Build upload module · Connect to FastAPI\nDisplay prediction results · Basic UI"),
    ("Tue, 24 March", "Human-in-the-Loop review UI · Progress bar\nTop-5 alternatives · Export · Testing"),
    ("Wed, 25 March — AM", "Polish UI · Bug fixes\nPrepare demo · Presentation rehearsal"),
    ("Wed, 25 March — PM  🏆", "Final Presentation & Live Demo"),
]
y = 1.7
for day, desc in timeline:
    # dot
    dot_color = PURPLE if "🏆" in day else TEAL
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(0.85), Inches(y + 0.08), Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()
    # line
    if timeline.index((day, desc)) < len(timeline) - 1:
        line = add_rect(s, 0.925, y + 0.28, 0.03, 0.7, fill_color=BORDER)
    add_textbox(s, 1.2, y, 5.5, 0.25, day,
                font_size=13, bold=True, color=TEAL_DARK)
    add_textbox(s, 1.2, y + 0.28, 5.5, 0.6, desc.replace("\n", "  ·  "),
                font_size=12, color=TEXT_SEC)
    y += 0.88

# Team roles (right)
add_textbox(s, 7.5, 1.7, 5, 0.35, "Suggested Team Roles",
            font_size=16, bold=True, color=TEXT_PRI)

roles = [
    ("🎨", "Frontend Developer(s)", "React/Vue components, UI, styling"),
    ("⚙️", "API Integration Lead", "FastAPI connection, data flow, errors"),
    ("🩺", "Clinical / Domain Expert", "MedDRA understanding, test cases, UX"),
    ("🏆", "Demo & Presentation Lead", "Final demo flow, slides, presentation"),
]
y = 2.2
for icon, title, desc in roles:
    add_rect(s, 7.5, y, 5.2, 0.75, fill_color=WHITE, border_color=BORDER,
             border_width=Pt(1.5))
    add_textbox(s, 7.7, y + 0.1, 0.5, 0.4, icon, font_size=22)
    add_textbox(s, 8.3, y + 0.1, 4.2, 0.3, title,
                font_size=14, bold=True, color=TEXT_PRI)
    add_textbox(s, 8.3, y + 0.42, 4.2, 0.3, desc,
                font_size=11, color=TEXT_DIM)
    y += 0.88

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 10 – LEARNING OUTCOMES & CTA
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1),
                            Inches(4), Inches(4))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0xE6, 0xFB, 0xF7)
circle.line.fill.background()

add_badge(s, 0.7, 0.5, "LEARNING OUTCOMES", TEAL_LIGHT, TEAL_DARK)
add_textbox(s, 0.7, 1.0, 10, 0.6, "What You'll Gain",
            font_size=32, bold=True, color=TEXT_PRI)

# Three skill columns
skill_cols = [
    ("🏥 Clinical Knowledge", TEAL_DARK, TEAL,
     ["MedDRA terminology & hierarchy",
      "Adverse event coding workflow",
      "Pharmacovigilance fundamentals",
      "Clinical trial safety reporting",
      "Human-in-the-Loop AI in medicine"]),
    ("💻 Technical Skills", PURPLE, PURPLE,
     ["REST API integration (FastAPI)",
      "Modern frontend development",
      "LLM pipeline architecture",
      "RAG system understanding",
      "Data upload/export workflows"]),
    ("🎯 Professional Experience", BLUE, BLUE,
     ["Agile sprint workflow",
      "Team collaboration & delivery",
      "Product demo & pitching",
      "Real-world Digital Health project",
      "Research-to-production thinking"]),
]
x = 0.7
for title, title_color, border_c, items in skill_cols:
    add_rect(s, x, 1.75, 3.8, 3.2, fill_color=WHITE,
             border_color=BORDER, border_width=Pt(1.5))
    # color top bar
    add_rect(s, x, 1.75, 3.8, 0.06, fill_color=border_c)
    add_textbox(s, x + 0.2, 1.9, 3.4, 0.35, title,
                font_size=13, bold=True, color=title_color)
    y_item = 2.35
    for item in items:
        add_textbox(s, x + 0.3, y_item, 3.3, 0.3, f"→  {item}",
                    font_size=12, color=TEXT_SEC)
        y_item += 0.35
    x += 4.1

# CTA Box
add_rect(s, 0.7, 5.2, 12, 1.9, fill_color=TEAL_LIGHT,
         border_color=TEAL, border_width=Pt(1.5))
add_textbox(s, 0.7, 5.35, 12, 0.5,
            "Ready to Build the Future of Drug Safety?",
            font_size=24, bold=True, color=TEAL_DARK, alignment=PP_ALIGN.CENTER)

cta_items = ["🧬 Real clinical data pipeline",
             "🤖 State-of-the-art LLM technology",
             "🏥 Direct impact on pharmacovigilance",
             "🚀 Research paper → Production tool"]
add_textbox(s, 0.7, 5.9, 12, 0.4,
            "    ·    ".join(cta_items),
            font_size=13, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)

add_textbox(s, 0.7, 6.45, 12, 0.35,
    "Supervisors: Isabella Catharina Wiest & Naghme Dashti  |  "
    "Else Kröner Fresenius Center for Digital Health, TU Dresden",
    font_size=12, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_footer(s, "EKFZ Digital Health  ·  TU Dresden  ·  Spring School March 2026")


# ═════════════════════════════════════════════════════════════
#  SAVE
# ═════════════════════════════════════════════════════════════
output = "AE_Coding_PitchDeck_Spring2026.pptx"
prs.save(output)
print(f"✅ Saved: {output}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Size: 16:9 widescreen")
print(f"\n📌 Don't forget to:")
print(f"   1. Replace logo placeholders on Slide 1")
print(f"   2. Add prototype screenshots on Slide 6")
print(f"   3. Review fonts (Calibri used throughout)")
