# ============================================================
#  Tumor Documentation Pitch Deck – Clinicum Digitale Spring 2026
#  Run: pip install python-pptx   then   python tumor_pitch.py
# ============================================================

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colours ──────────────────────────────────────────────────
TEAL       = RGBColor(0x0D, 0x94, 0x88)
TEAL_DARK  = RGBColor(0x0F, 0x76, 0x6E)
TEAL_LIGHT = RGBColor(0xCC, 0xFB, 0xF1)
PURPLE     = RGBColor(0x63, 0x66, 0xF1)
PURPLE_LT  = RGBColor(0xEE, 0xF2, 0xFF)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_LT   = RGBColor(0xFF, 0xFB, 0xEB)
RED        = RGBColor(0xEF, 0x44, 0x44)
RED_LT     = RGBColor(0xFE, 0xF2, 0xF2)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG_SUBTLE  = RGBColor(0xF1, 0xF5, 0xF9)
BG_LIGHT   = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_PRI   = RGBColor(0x0F, 0x17, 0x2A)
TEXT_SEC   = RGBColor(0x47, 0x55, 0x69)
TEXT_DIM   = RGBColor(0x94, 0xA3, 0xB8)
BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
BLUE       = RGBColor(0x25, 0x63, 0xEB)
BLUE_LT    = RGBColor(0xEF, 0xF6, 0xFF)
ORANGE     = RGBColor(0xEA, 0x58, 0x0C)
ORANGE_LT  = RGBColor(0xFF, 0xF7, 0xED)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)
GREEN_LT   = RGBColor(0xF0, 0xFD, 0xF4)

# ── Presentation Setup ───────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper functions ─────────────────────────────────────────
def add_blank_slide():
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text,
                font_size=14, bold=False, color=TEXT_PRI,
                alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_textbox(slide, left, top, width, height, runs,
                     alignment=PP_ALIGN.LEFT, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing * 14)
    for i, r in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
        else:
            run = p.add_run()
        run.text = r.get("text", "")
        run.font.size = Pt(r.get("size", 14))
        run.font.bold = r.get("bold", False)
        run.font.color.rgb = r.get("color", TEXT_PRI)
        run.font.name = r.get("font_name", "Calibri")
    return txBox

def add_rect(slide, left, top, width, height, fill_color=BG_SUBTLE,
             border_color=None, border_width=Pt(1), corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=13, color=TEXT_SEC, bold=False,
                          alignment=PP_ALIGN.LEFT, line_spacing=1.4,
                          font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
    return txBox

def add_badge(slide, left, top, text, bg_color=TEAL_LIGHT, text_color=TEAL_DARK):
    w, h = 3.5, 0.4
    shape = add_rect(slide, left, top, w, h, fill_color=bg_color)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(10)
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.color.rgb = text_color
    shape.text_frame.paragraphs[0].font.name = "Calibri"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_footer(slide, text="EKFZ Digital Health  ·  TU Dresden  ·  March 2026"):
    add_textbox(slide, 0.6, 7.0, 5, 0.3, text,
                font_size=9, color=TEXT_DIM)

def add_card_box(slide, left, top, width, height, title, body_lines,
                 icon="", border_color=BORDER):
    add_rect(slide, left, top, width, height,
             fill_color=WHITE, border_color=border_color, border_width=Pt(1.5))
    y = top + 0.2
    if icon:
        add_textbox(slide, left + 0.25, y, 0.5, 0.4, icon, font_size=22)
        y += 0.35
    add_textbox(slide, left + 0.25, y, width - 0.5, 0.35, title,
                font_size=15, bold=True, color=TEXT_PRI)
    y += 0.4
    add_multiline_textbox(slide, left + 0.25, y, width - 0.5,
                          height - (y - top) - 0.2, body_lines,
                          font_size=12, color=TEXT_SEC)

def add_pipeline_box(slide, left, top, w, h, num, title, desc_lines, active=True):
    bc = ORANGE if active else BORDER
    add_rect(slide, left, top, w, h, fill_color=WHITE,
             border_color=bc, border_width=Pt(2))
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left + w/2 - 0.2), Inches(top + 0.2),
        Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = str(num)
    circle.text_frame.paragraphs[0].font.size = Pt(12)
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_textbox(slide, left + 0.15, top + 0.7, w - 0.3, 0.35, title,
                font_size=14, bold=True, color=TEXT_PRI,
                alignment=PP_ALIGN.CENTER)
    add_multiline_textbox(slide, left + 0.15, top + 1.05, w - 0.3, h - 1.25,
                          desc_lines, font_size=11, color=TEXT_SEC,
                          alignment=PP_ALIGN.CENTER)

def add_arrow_text(slide, left, top, text="→"):
    add_textbox(slide, left, top, 0.4, 0.5, text,
                font_size=24, bold=True, color=ORANGE,
                alignment=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════
#  SLIDE 1 – TITLE
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()

# Decorative circle
circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1),
                            Inches(4), Inches(4))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0xFF, 0xF7, 0xED)
circle.line.fill.background()

add_badge(s, 0.7, 0.7,
          "CLINICUM DIGITALE  ·  SPRING 2026  ·  PROJECT X",
          ORANGE_LT, ORANGE)

add_rich_textbox(s, 0.7, 1.4, 10, 1.5, [
    {"text": "LLM-Powered ", "size": 42, "bold": True, "color": TEXT_PRI},
    {"text": "Tumor\n", "size": 42, "bold": True, "color": ORANGE},
    {"text": "Documentation Tool", "size": 42, "bold": True, "color": TEXT_PRI},
])

add_textbox(s, 0.7, 3.2, 9, 0.8,
    "Build a web tool that extracts structured oncological information from "
    "clinical documents — powered by LLM-based Retrieval-Augmented Generation "
    "and ICD-O coding.",
    font_size=17, color=TEXT_SEC)

# Supervisor box
add_rect(s, 0.7, 4.5, 7, 0.9, fill_color=BG_SUBTLE)
add_textbox(s, 1.0, 4.6, 6.5, 0.3,
            "Isabella Catharina Wiest & Naghme Dashti",
            font_size=16, bold=True, color=TEXT_PRI)
add_textbox(s, 1.0, 4.95, 6.5, 0.3,
            "Else Kröner Fresenius Center for Digital Health  ·  TU Dresden",
            font_size=12, color=TEXT_DIM)

# Logo placeholders
add_rect(s, 9.5, 6.2, 1.2, 0.8, fill_color=BG_SUBTLE, border_color=BORDER)
add_textbox(s, 9.5, 6.3, 1.2, 0.6, "EKFZ\nLogo", font_size=9,
            color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_rect(s, 10.9, 6.2, 1.2, 0.8, fill_color=BG_SUBTLE, border_color=BORDER)
add_textbox(s, 10.9, 6.3, 1.2, 0.6, "TUD\nLogo", font_size=9,
            color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_rect(s, 12.3, 6.2, 0.8, 0.8, fill_color=BG_SUBTLE, border_color=BORDER)
add_textbox(s, 12.3, 6.3, 0.8, 0.6, "CGC\nLogo", font_size=8,
            color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 2 – THE PROBLEM
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
add_badge(s, 0.7, 0.5, "THE CHALLENGE", PURPLE_LT, PURPLE)
add_textbox(s, 0.7, 1.0, 10, 0.6, "The Tumor Documentation Burden",
            font_size=32, bold=True, color=TEXT_PRI)

# Problem points (left)
problems = [
    ("🏥  Oncology generates massive unstructured data",
     "Physician letters, pathology reports, follow-up notes, imaging summaries — "
     "all as free-text documents."),
    ("📋  Manual documentation is a bottleneck",
     "Transferring tumor info into structured registries is time-consuming, "
     "error-prone, and burdens clinical staff."),
    ("🔗  Fragmented patient records",
     "Relevant tumor data is scattered across multiple documents per patient — "
     "morphology, topography, staging, grading, therapy."),
    ("⚠️  Incomplete documentation affects care",
     "Missing or delayed structured data impacts therapy decisions, "
     "interdisciplinary collaboration, and tumor board outcomes."),
]

y = 1.75
for title, desc in problems:
    add_rect(s, 0.7, y, 6.0, 0.72, fill_color=WHITE, border_color=BORDER)
    add_textbox(s, 0.9, y + 0.05, 5.6, 0.3, title,
                font_size=13, bold=True, color=TEXT_PRI)
    add_textbox(s, 0.9, y + 0.35, 5.6, 0.35, desc,
                font_size=11, color=TEXT_SEC)
    y += 0.85

# Right side — Current vs Vision
add_rect(s, 7.2, 1.75, 5.4, 1.2, fill_color=RED_LT, border_color=RED,
         border_width=Pt(1.5))
add_textbox(s, 7.4, 1.85, 5, 0.3, "❌  Current Workflow",
            font_size=14, bold=True, color=RED)
add_textbox(s, 7.4, 2.2, 5, 0.3,
            "📄 Multiple Docs  →  👨‍⚕️ Manual Reading  →  📋 Registry Entry",
            font_size=12, color=TEXT_SEC)
add_textbox(s, 7.4, 2.55, 5, 0.3,
            "⏳ 30+ min per patient  ·  ❌ Error-prone  ·  📉 Incomplete",
            font_size=11, bold=True, color=RED)

add_rect(s, 7.2, 3.15, 5.4, 1.2, fill_color=ORANGE_LT, border_color=ORANGE,
         border_width=Pt(1.5))
add_textbox(s, 7.4, 3.25, 5, 0.3, "✅  Our Vision",
            font_size=14, bold=True, color=ORANGE)
add_textbox(s, 7.4, 3.6, 5, 0.3,
            "📄 Upload Docs  →  🤖 LLM Pipeline  →  ✓ Structured Output",
            font_size=12, color=TEXT_SEC)
add_textbox(s, 7.4, 3.95, 5, 0.3,
            "⚡ Automated  ·  🎯 Structured ICD-O  ·  📈 Complete",
            font_size=11, bold=True, color=ORANGE)

# Statistics box
add_rect(s, 7.2, 4.6, 5.4, 1.5, fill_color=BG_SUBTLE)
add_textbox(s, 7.4, 4.7, 5.0, 0.3, "📊  Real-World Numbers",
            font_size=14, bold=True, color=TEXT_PRI)
add_multiline_textbox(s, 7.4, 5.05, 5.0, 1.0, [
    "• 77 OCR documents across 50 patients",
    "• Multiple document types per patient",
    "• 63 tumor-level ground truth labels",
    "• Some patients have multiple tumors",
], font_size=12, color=TEXT_SEC)

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 3 – ICD-O CODING BACKGROUND
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
add_badge(s, 0.7, 0.5, "BACKGROUND", ORANGE_LT, ORANGE)
add_textbox(s, 0.7, 1.0, 10, 0.6, "ICD-O: The Oncology Coding Standard",
            font_size=32, bold=True, color=TEXT_PRI)
add_textbox(s, 0.7, 1.5, 10, 0.35,
            "International Classification of Diseases for Oncology — "
            "dual-axis coding: WHERE (Topography) + WHAT (Morphology)",
            font_size=14, color=TEXT_SEC)

# Left: Two coding axes
# Topography box
add_rect(s, 0.7, 2.2, 5.5, 2.0, fill_color=WHITE, border_color=BLUE,
         border_width=Pt(2))
add_rect(s, 0.7, 2.2, 5.5, 0.06, fill_color=BLUE)
add_textbox(s, 0.9, 2.35, 5.1, 0.35, "📍  Topography — WHERE is the tumor?",
            font_size=16, bold=True, color=BLUE)
add_multiline_textbox(s, 0.9, 2.75, 5.1, 1.3, [
    "• 1,338 terms in ICD-O-3 dictionary",
    "• Format: C00–C80 (e.g., C61 = Prostate, C50 = Breast)",
    "• Specifies anatomical site of primary tumor",
    "• Examples: C73.9 Thyroid, C56.9 Ovary, C18.0 Cecum",
], font_size=13, color=TEXT_SEC)

# Morphology box
add_rect(s, 0.7, 4.45, 5.5, 2.0, fill_color=WHITE, border_color=PURPLE,
         border_width=Pt(2))
add_rect(s, 0.7, 4.45, 5.5, 0.06, fill_color=PURPLE)
add_textbox(s, 0.9, 4.6, 5.1, 0.35, "🔬  Morphology — WHAT type of tumor?",
            font_size=16, bold=True, color=PURPLE)
add_multiline_textbox(s, 0.9, 5.0, 5.1, 1.3, [
    "• 2,903 terms in ICD-O-3 dictionary",
    "• Format: XXXX/B (e.g., 8140/3 = Adenocarcinoma)",
    "• Describes histology + behavior (/0 benign … /3 malignant)",
    "• Examples: 8500/3 Ductal carcinoma, 8720/3 Melanoma",
], font_size=13, color=TEXT_SEC)

# Right: Document types + extraction
add_textbox(s, 6.8, 2.1, 5.5, 0.35, "Clinical Document Types",
            font_size=16, bold=True, color=ORANGE)

doc_types = [
    ("📝", "Arztbrief", "Physician letter with diagnosis,\ntherapy plan, clinical findings"),
    ("🔬", "Pathologiebefund", "Histological analysis, grading,\nmorphology, resection margins"),
    ("📊", "Verlaufsdoku", "Follow-up documentation,\nstaging, response assessment"),
    ("🖼", "Befundbericht", "Imaging/diagnostic reports,\ntumor size, metastasis status"),
]

y = 2.55
for icon, title, desc in doc_types:
    add_rect(s, 6.8, y, 5.5, 0.82, fill_color=WHITE, border_color=BORDER)
    add_textbox(s, 7.0, y + 0.1, 0.4, 0.35, icon, font_size=20)
    add_textbox(s, 7.5, y + 0.08, 2.0, 0.3, title,
                font_size=14, bold=True, color=TEXT_PRI)
    add_textbox(s, 7.5, y + 0.38, 4.5, 0.4,
                desc.replace("\n", " "),
                font_size=11, color=TEXT_SEC)
    y += 0.92

# Bottom highlight
add_rect(s, 6.8, 6.25, 5.5, 0.55, fill_color=ORANGE_LT)
add_rect(s, 6.8, 6.25, 0.06, 0.55, fill_color=ORANGE)
add_textbox(s, 7.05, 6.3, 5.1, 0.45,
    "Each document type requires a specific prompt\n"
    "to extract the right oncological information.",
    font_size=12, bold=True, color=ORANGE)

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 4 – PIPELINE OVERVIEW
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
add_badge(s, 0.7, 0.5, "OUR APPROACH", ORANGE_LT, ORANGE)
add_textbox(s, 0.7, 1.0, 10, 0.6, "Multi-Stage LLM Pipeline",
            font_size=32, bold=True, color=TEXT_PRI)
add_textbox(s, 0.7, 1.55, 10, 0.3,
            "Four stages: Document Ingestion → LLM Summarization → RAG Retrieval → ICD-O Assignment",
            font_size=14, color=TEXT_SEC)

# Four pipeline stages
stages = [
    (1, "Document\nIngestion",
     ["OCR-processed clinical docs",
      "Multiple docs per patient",
      "Merge into unified text",
      "Clean & chunk text"]),
    (2, "LLM\nSummarization",
     ["GPT-OSS-120B on Pluto",
      "Tumor-focused extraction",
      "Map-reduce chunking",
      "7–10 sentence summary"]),
    (3, "RAG\nRetrieval",
     ["BGE-M3 embeddings",
      "Cosine similarity search",
      "Top-K Morph candidates",
      "Top-K Topo candidates"]),
    (4, "ICD-O Code\nAssignment",
     ["LLM reasoning over",
      "summary + candidates",
      "Structured JSON output",
      "Morph + Topo + Reason"]),
]

x = 0.5
for num, title, desc in stages:
    add_pipeline_box(s, x, 2.2, 2.8, 2.8, num, title.replace("\n", " "), desc)
    if num < 4:
        add_arrow_text(s, x + 2.85, 3.4)
    x += 3.15

# Bottom boxes
add_rect(s, 0.7, 5.4, 12, 0.8, fill_color=ORANGE_LT)
add_rect(s, 0.7, 5.4, 0.06, 0.8, fill_color=ORANGE)
add_textbox(s, 1.0, 5.5, 11.5, 0.6,
    "💡  The entire LLM pipeline (Stages 1–4) already exists as a FastAPI backend — "
    "your mission is to build the web interface for document upload, type assignment, "
    "and result display!",
    font_size=15, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

# Tech labels
tech = ["🧠 LLM: GPT-OSS-120B / LLaMA-3.3-70B",
        "📐 Embeddings: BGE-M3 (1024-dim)",
        "📖 ICD-O-3: 2903 Morph + 1338 Topo",
        "⚡ API: FastAPI backend"]
x = 0.7
for t in tech:
    add_rect(s, x, 6.4, 3.0, 0.45, fill_color=BG_SUBTLE)
    add_textbox(s, x, 6.45, 3.0, 0.35, t,
                font_size=11, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)
    x += 3.15

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 5 – CURRENT STATUS & RESULTS
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
add_badge(s, 0.7, 0.5, "CURRENT STATUS", PURPLE_LT, PURPLE)
add_textbox(s, 0.7, 1.0, 10, 0.6,
            "Where We Are — Research Results",
            font_size=32, bold=True, color=TEXT_PRI)
add_textbox(s, 0.7, 1.55, 10, 0.3,
            "Baseline evaluation on 50 patients with patient-level ANY-match strategy",
            font_size=14, color=TEXT_SEC)

# Current stats
stats = [
    ("50", "Patients"),
    ("77", "OCR Documents"),
    ("63", "GT Tumor Labels"),
    ("4,241", "ICD-O Dictionary\nTerms"),
]
x = 0.7
for value, label in stats:
    add_rect(s, x, 2.1, 2.8, 1.0, fill_color=WHITE, border_color=BORDER,
             border_width=Pt(1.5))
    add_textbox(s, x, 2.15, 2.8, 0.45, value,
                font_size=28, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x, 2.65, 2.8, 0.4, label,
                font_size=11, bold=True, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    x += 3.05

# Left: Key findings
add_textbox(s, 0.7, 3.4, 6, 0.35, "Key Research Findings",
            font_size=16, bold=True, color=TEXT_PRI)

findings = [
    ("✅", "LLM reasoning is strong",
     "When correct code is in candidates: 81% Morph / 78% Topo accuracy",
     GREEN_LT, GREEN),
    ("⚠️", "Retrieval is the bottleneck",
     "Correct code in Top-K candidates only 23% (Morph) / 19% (Topo)",
     AMBER_LT, RGBColor(0xB4, 0x53, 0x09)),
    ("📊", "NOS code bias observed",
     'Generic "NOS" codes over-represented; specific codes under-retrieved',
     PURPLE_LT, PURPLE),
    ("🔍", "Strict matching underestimates",
     "C61 vs C61.9 fails strict match — relaxed Topo accuracy: 21%",
     BLUE_LT, BLUE),
]

y = 3.85
for icon, title, desc, bg, tc in findings:
    add_rect(s, 0.7, y, 6.0, 0.72, fill_color=bg, border_color=tc,
             border_width=Pt(1))
    add_textbox(s, 0.9, y + 0.05, 5.6, 0.3,
                f"{icon}  {title}", font_size=13, bold=True, color=tc)
    add_textbox(s, 0.9, y + 0.37, 5.6, 0.3, desc,
                font_size=11, color=TEXT_SEC)
    y += 0.82

# Right: What this means for the sprint
add_rect(s, 7.2, 3.4, 5.4, 3.45, fill_color=WHITE, border_color=BORDER,
         border_width=Pt(1.5))
add_textbox(s, 7.4, 3.5, 5.0, 0.35,
            "What This Means for Your Sprint",
            font_size=16, bold=True, color=TEXT_PRI)
add_multiline_textbox(s, 7.4, 3.95, 5.0, 2.8, [
    "The pipeline works — the LLM is capable of making",
    "good coding decisions when given the right candidates.",
    "",
    "What's missing is a user-facing tool that lets",
    "clinicians:",
    "",
    "  📤  Upload & organize their documents",
    "  🏷  Assign document types for prompt selection",
    "  ▶️  Trigger the pipeline via the API",
    "  📋  View structured results clearly",
    "  ✏️  Review & correct predictions",
    "",
    "→ That's exactly what you'll build!",
], font_size=13, color=TEXT_SEC)

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 6 – WHAT EXISTS (CURRENT WORK)
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
add_badge(s, 0.7, 0.5, "WHAT EXISTS", ORANGE_LT, ORANGE)
add_textbox(s, 0.7, 1.0, 10, 0.6, "Current Pipeline Components",
            font_size=32, bold=True, color=TEXT_PRI)
add_textbox(s, 0.7, 1.55, 10, 0.3,
            "Research pipeline runs as scripts — no web interface yet",
            font_size=14, color=TEXT_SEC)

# Three existing components
add_card_box(s, 0.7, 2.1, 3.7, 3.5,
    "📄  OCR & Summarization",
    ["Task 2.1 — Complete ✅",
     "",
     "• PaddleOCR → 77 text files",
     "• Merged to 50 patient docs",
     "• GPT-OSS-120B summarization",
     "• Map-reduce chunking strategy",
     "• Robust multi-stage fallback",
     "• Output: structured JSON + CSV"],
    border_color=GREEN)

add_card_box(s, 4.7, 2.1, 3.7, 3.5,
    "📐  Embeddings & Retrieval",
    ["Task 2.2 + 2.3 — Complete ✅",
     "",
     "• BGE-M3 embeddings (1024-dim)",
     "• 50 summary embeddings",
     "• 2,903 morphology embeddings",
     "• 1,338 topography embeddings",
     "• Cosine similarity Top-K retrieval",
     "• Hybrid semantic + lexical (v3)"],
    border_color=GREEN)

add_card_box(s, 8.7, 2.1, 3.9, 3.5,
    "🧠  RAG ICD-O Prediction",
    ["Task 2.3 — In Progress 🔄",
     "",
     "• LLM reads summary + candidates",
     "• Picks best Morph + Topo code",
     "• Outputs structured JSON",
     "• 3 pipeline versions (v1→v3)",
     "• v3: organ-based filtering",
     "• FastAPI wrapper available"],
    border_color=ORANGE)

# Screenshot placeholder
add_rect(s, 0.7, 5.85, 5.5, 0.6, fill_color=BG_SUBTLE, border_color=BORDER)
add_textbox(s, 0.7, 5.9, 5.5, 0.5,
            "📸 Placeholder — Patient Summary Example Output",
            font_size=10, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_rect(s, 6.5, 5.85, 5.5, 0.6, fill_color=BG_SUBTLE, border_color=BORDER)
add_textbox(s, 6.5, 5.9, 5.5, 0.5,
            "📸 Placeholder — ICD-O Prediction JSON Example",
            font_size=10, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# CTA
add_rect(s, 0.7, 6.6, 12, 0.5, fill_color=PURPLE_LT)
add_textbox(s, 0.7, 6.65, 12, 0.4,
    "🚀  The pipeline runs as Python scripts — now let's make it a usable web tool!",
    font_size=15, bold=True, color=PURPLE, alignment=PP_ALIGN.CENTER)

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 7 – YOUR MISSION
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
add_badge(s, 0.7, 0.5, "YOUR MISSION", PURPLE_LT, PURPLE)
add_textbox(s, 0.7, 1.0, 10, 0.6, "What You Will Build 🛠️",
            font_size=32, bold=True, color=TEXT_PRI)

# Left: Feature list
features = [
    ("📤", "Drag-and-Drop Document Upload",
     "Upload multiple medical documents (PDF, text) with preview"),
    ("🏷", "Document Type Assignment",
     "Assign each document: Arztbrief, Pathologiebefund, Verlaufsdoku, etc."),
    ("⚡", "Pipeline Trigger via FastAPI",
     "Start LLM pipeline with document-type-specific prompts per document"),
    ("📊", "Structured Results Display",
     "Show extracted tumor info: ICD-O Topography, Morphology, Grading, Stage"),
    ("✏️", "Review & Edit Interface",
     "Clinicians can review, correct, and confirm extracted information"),
    ("💾", "Export Structured Data",
     "Download results as JSON / CSV for registry integration"),
]
y = 1.7
for icon, title, desc in features:
    add_rect(s, 0.7, y, 6.8, 0.65, fill_color=WHITE, border_color=BORDER)
    add_textbox(s, 0.85, y + 0.05, 0.4, 0.3, icon, font_size=18)
    add_textbox(s, 1.35, y + 0.05, 3, 0.28, title,
                font_size=13, bold=True, color=TEXT_PRI)
    add_textbox(s, 1.35, y + 0.33, 5.8, 0.28, desc,
                font_size=11, color=TEXT_SEC)
    y += 0.73

# Right: Build vs Provided
add_rect(s, 8.0, 1.7, 4.6, 2.2, fill_color=WHITE,
         border_color=ORANGE, border_width=Pt(2))
add_rect(s, 8.15, 1.85, 1.7, 0.3, fill_color=ORANGE_LT)
add_textbox(s, 8.15, 1.85, 1.7, 0.3, "🧑‍💻 YOU BUILD",
            font_size=9, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
add_multiline_textbox(s, 8.2, 2.25, 4.2, 1.5, [
    "• Document upload & preview UI",
    "• Document type selector",
    "• API client for pipeline trigger",
    "• Results display & visualization",
    "• Review/edit workflow",
    "• Export functionality",
], font_size=12, color=TEXT_SEC)

add_rect(s, 8.0, 4.05, 4.6, 1.5, fill_color=WHITE,
         border_color=PURPLE, border_width=Pt(2))
add_rect(s, 8.15, 4.2, 1.7, 0.3, fill_color=PURPLE_LT)
add_textbox(s, 8.15, 4.2, 1.7, 0.3, "🎁 PROVIDED",
            font_size=9, bold=True, color=PURPLE, alignment=PP_ALIGN.CENTER)
add_multiline_textbox(s, 8.2, 4.6, 4.2, 0.85, [
    "• FastAPI backend (full pipeline)",
    "• LLM summarization + RAG retrieval",
    "• ICD-O-3 dictionary embeddings",
    "• Document-type-specific prompts",
], font_size=12, color=TEXT_SEC)

add_rect(s, 8.0, 5.7, 4.6, 0.8, fill_color=RED_LT,
         border_color=RED, border_width=Pt(1.5))
add_textbox(s, 8.2, 5.75, 4.2, 0.3, "❌ Out of Scope",
            font_size=12, bold=True, color=RED)
add_textbox(s, 8.2, 6.1, 4.2, 0.35,
            "No NLP/OCR/AI dev, no model training,\nno medical validation",
            font_size=11, color=TEXT_SEC)

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 8 – ARCHITECTURE
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
add_badge(s, 0.7, 0.5, "TECHNICAL ARCHITECTURE", ORANGE_LT, ORANGE)
add_textbox(s, 0.7, 1.0, 10, 0.6, "System Overview",
            font_size=32, bold=True, color=TEXT_PRI)

# Frontend box
add_rect(s, 1.0, 1.8, 11.3, 1.8, fill_color=WHITE,
         border_color=ORANGE, border_width=Pt(2))
add_rect(s, 1.15, 1.95, 2.0, 0.3, fill_color=ORANGE_LT)
add_textbox(s, 1.15, 1.95, 2.0, 0.3, "🧑‍💻 FRONTEND — You Build",
            font_size=9, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

fe_modules = [
    ("📤", "Upload Module", "Drag & drop + type assign"),
    ("▶️", "Pipeline Control", "Trigger & progress tracking"),
    ("📊", "Results View", "Structured tumor data"),
    ("✏️", "Review Panel", "Edit & confirm results"),
]
x = 1.3
for icon, title, desc in fe_modules:
    add_rect(s, x, 2.4, 2.5, 0.95, fill_color=BG_SUBTLE)
    add_textbox(s, x + 0.1, 2.45, 2.3, 0.35, f"{icon}  {title}",
                font_size=13, bold=True, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + 0.1, 2.8, 2.3, 0.35, desc,
                font_size=11, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    x += 2.7

# Arrow
add_textbox(s, 5.5, 3.7, 2.5, 0.5, "⇅  REST API  ⇅",
            font_size=18, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

# Backend box
add_rect(s, 1.0, 4.3, 11.3, 1.8, fill_color=WHITE,
         border_color=PURPLE, border_width=Pt(2))
add_rect(s, 1.15, 4.45, 2.2, 0.3, fill_color=PURPLE_LT)
add_textbox(s, 1.15, 4.45, 2.2, 0.3, "🎁 BACKEND — Provided",
            font_size=9, bold=True, color=PURPLE, alignment=PP_ALIGN.CENTER)

be_modules = [
    ("📄", "Doc Processing", "OCR + merging + cleaning"),
    ("🧠", "LLM Summarization", "GPT-OSS-120B / LLaMA"),
    ("📐", "RAG Retrieval", "BGE-M3 cosine Top-K"),
    ("🏷", "ICD-O Assignment", "Morph + Topo + Reason"),
]
x = 1.3
for icon, title, desc in be_modules:
    add_rect(s, x, 4.9, 2.5, 0.95, fill_color=BG_SUBTLE)
    add_textbox(s, x + 0.1, 4.95, 2.3, 0.35, f"{icon}  {title}",
                font_size=13, bold=True, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + 0.1, 5.3, 2.3, 0.35, desc,
                font_size=11, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    x += 2.7

# Tech stack
tech = ["🖥 Frontend: React / Vue / Streamlit",
        "⚡ API: FastAPI (provided)",
        "🧠 LLMs: GPT-OSS-120B / LLaMA-3.3",
        "📐 Embeddings: BGE-M3"]
x = 1.0
for t in tech:
    add_rect(s, x, 6.4, 2.8, 0.45, fill_color=BG_SUBTLE)
    add_textbox(s, x, 6.45, 2.8, 0.35, t,
                font_size=11, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)
    x += 3.0

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 9 – TIMELINE & TEAM
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
add_badge(s, 0.7, 0.5, "SPRINT PLAN", AMBER_LT, RGBColor(0xB4, 0x53, 0x09))
add_textbox(s, 0.7, 1.0, 10, 0.6, "Project Timeline & Team Roles",
            font_size=32, bold=True, color=TEXT_PRI)

# Timeline (left)
timeline = [
    ("Fri, 20 March — PM",
     "Kick-off · Understand pipeline & API docs · Set up dev environment · Team roles"),
    ("Mon, 23 March",
     "Build upload module with drag & drop · Document type selector · API connection"),
    ("Tue, 24 March",
     "Results display · Review/edit panel · Progress tracking · Export · Testing"),
    ("Wed, 25 March — AM",
     "Polish UI · Bug fixes · Prepare demo · Presentation rehearsal"),
    ("Wed, 25 March — PM  🏆",
     "Final Presentation & Live Demo"),
]
y = 1.7
for i, (day, desc) in enumerate(timeline):
    dot_color = PURPLE if "🏆" in day else ORANGE
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(0.85), Inches(y + 0.08), Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()
    if i < len(timeline) - 1:
        add_rect(s, 0.925, y + 0.28, 0.03, 0.7, fill_color=BORDER)
    add_textbox(s, 1.2, y, 5.5, 0.25, day,
                font_size=13, bold=True, color=ORANGE)
    add_textbox(s, 1.2, y + 0.28, 5.5, 0.6, desc,
                font_size=12, color=TEXT_SEC)
    y += 0.88

# Team roles (right)
add_textbox(s, 7.5, 1.7, 5, 0.35, "Suggested Team Roles",
            font_size=16, bold=True, color=TEXT_PRI)

roles = [
    ("🎨", "Frontend Developer(s)", "Upload UI, results view, responsive design"),
    ("⚙️", "API Integration Lead", "FastAPI connection, pipeline trigger, data flow"),
    ("🩺", "Clinical / Domain Expert", "ICD-O understanding, document types, UX review"),
    ("🏆", "Demo & Presentation Lead", "Final demo flow, slides, live presentation"),
]
y = 2.2
for icon, title, desc in roles:
    add_rect(s, 7.5, y, 5.2, 0.75, fill_color=WHITE, border_color=BORDER,
             border_width=Pt(1.5))
    add_textbox(s, 7.7, y + 0.1, 0.5, 0.4, icon, font_size=22)
    add_textbox(s, 8.3, y + 0.1, 4.2, 0.3, title,
                font_size=14, bold=True, color=TEXT_PRI)
    add_textbox(s, 8.3, y + 0.42, 4.2, 0.3, desc,
                font_size=11, color=TEXT_DIM)
    y += 0.88

add_footer(s)


# ═════════════════════════════════════════════════════════════
#  SLIDE 10 – LEARNING OUTCOMES & CTA
# ═════════════════════════════════════════════════════════════
s = add_blank_slide()
circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1),
                            Inches(4), Inches(4))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0xFF, 0xF7, 0xED)
circle.line.fill.background()

add_badge(s, 0.7, 0.5, "LEARNING OUTCOMES", ORANGE_LT, ORANGE)
add_textbox(s, 0.7, 1.0, 10, 0.6, "What You'll Gain",
            font_size=32, bold=True, color=TEXT_PRI)

# Three columns
skill_cols = [
    ("🏥 Clinical Knowledge", ORANGE, ORANGE,
     ["ICD-O oncology coding system",
      "Clinical document types",
      "Tumor documentation workflow",
      "Structured vs unstructured data",
      "Interdisciplinary oncology processes"]),
    ("💻 Technical Skills", PURPLE, PURPLE,
     ["REST API integration (FastAPI)",
      "Modern frontend development",
      "LLM pipeline architecture",
      "RAG system understanding",
      "Document processing workflows"]),
    ("🎯 Professional Experience", BLUE, BLUE,
     ["Agile sprint workflow",
      "Team collaboration & delivery",
      "Product demo & pitching",
      "Real-world Digital Health project",
      "Research-to-production thinking"]),
]
x = 0.7
for title, title_color, border_c, items in skill_cols:
    add_rect(s, x, 1.75, 3.8, 3.2, fill_color=WHITE,
             border_color=BORDER, border_width=Pt(1.5))
    add_rect(s, x, 1.75, 3.8, 0.06, fill_color=border_c)
    add_textbox(s, x + 0.2, 1.9, 3.4, 0.35, title,
                font_size=13, bold=True, color=title_color)
    y_item = 2.35
    for item in items:
        add_textbox(s, x + 0.3, y_item, 3.3, 0.3, f"→  {item}",
                    font_size=12, color=TEXT_SEC)
        y_item += 0.35
    x += 4.1

# CTA Box
add_rect(s, 0.7, 5.2, 12, 1.9, fill_color=ORANGE_LT,
         border_color=ORANGE, border_width=Pt(1.5))
add_textbox(s, 0.7, 5.35, 12, 0.5,
            "Ready to Transform Oncology Documentation?",
            font_size=24, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

cta_items = ["🧬 Real clinical oncology data",
             "🤖 State-of-the-art LLM + RAG",
             "🏥 Direct impact on tumor registries",
             "🚀 Research pipeline → Production tool"]
add_textbox(s, 0.7, 5.9, 12, 0.4,
            "    ·    ".join(cta_items),
            font_size=13, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)

add_textbox(s, 0.7, 6.45, 12, 0.35,
    "Supervisors: Isabella Catharina Wiest & Naghme Dashti  |  "
    "Else Kröner Fresenius Center for Digital Health, TU Dresden",
    font_size=12, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_footer(s, "EKFZ Digital Health  ·  TU Dresden  ·  Spring School March 2026")


# ═════════════════════════════════════════════════════════════
#  SAVE
# ═════════════════════════════════════════════════════════════
output = "Tumor_Documentation_PitchDeck_Spring2026.pptx"
prs.save(output)
print(f"✅ Saved: {output}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Size: 16:9 widescreen")
print(f"\n📌 Don't forget to:")
print(f"   1. Replace logo placeholders on Slide 1")
print(f"   2. Add patient summary screenshot on Slide 6")
print(f"   3. Add ICD-O prediction JSON example on Slide 6")
